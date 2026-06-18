"""
app.py — Dashboard Complejo Térmico Mejillones v5
Fixes: dots de color en tabs, estado conexión sidebar,
eje X visible, títulos separados, programada visible
"""

import streamlit as st
from streamlit_autorefresh import st_autorefresh
import pandas as pd
import psycopg2
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import date, datetime, timedelta
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, white
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage

st.set_page_config(
    page_title="Complejo Térmico Mejillones",
    layout="wide",
    page_icon=None,
    initial_sidebar_state="expanded",
)

# Auto-refresh cada 60 minutos (3600000 ms) — mantiene la app despierta y datos frescos
st_autorefresh(interval=3_600_000, limit=None, key="autorefresh_horario")

# Paleta AES: púrpura → azul → cyan → verde (logo AES)
COLORES = {
    "ANG1": {"line":"#6D28D9","prog":"#C4B5FD","badge":"#EDE9FE","text":"#6D28D9"},
    "ANG2": {"line":"#2563EB","prog":"#93C5FD","badge":"#DBEAFE","text":"#2563EB"},
    "CCR1": {"line":"#0891B2","prog":"#67E8F9","badge":"#CFFAFE","text":"#0891B2"},
    "CCR2": {"line":"#16A34A","prog":"#86EFAC","badge":"#DCFCE7","text":"#16A34A"},
    "CMG":  {"line":"#6D28D9"},
}
LABELS = {"ANG1":"Angamos U1","ANG2":"Angamos U2","CCR1":"Cochrane U1","CCR2":"Cochrane U2"}

# Potencias máximas reales declaradas por AES Andes ante el CEN
PMAX = {"ANG1": 277.0, "ANG2": 280.0, "CCR1": 276.0, "CCR2": 276.0}

NOMBRES_NODO = {
    "CRUCERO_______220": "Crucero 220kV",
    "TARAPACA______220": "Tarapacá 220kV",
}

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Mono:wght@400;600&family=Inter:wght@300;400;500;600;700&display=swap');
:root{--bg:#F1F5F9;--surf:#FFFFFF;--surf2:#F8FAFC;--bord:#E2E8F0;--txt:#0F172A;--muted:#64748B;--accent:#2563EB;}
.stApp{background:var(--bg)!important;}
.block-container{padding:1.5rem 2rem 3rem;max-width:1400px;}
[data-testid="stSidebar"]{background:#1E293B!important;border-right:1px solid #334155!important;}
[data-testid="stSidebar"] *{color:#E2E8F0!important;}
[data-testid="stSidebar"] .status-box{background:#0F172A!important;border-color:#334155!important;}
[data-testid="stSidebar"] .stCheckbox label{color:#CBD5E1!important;}
[data-testid="stSidebar"] .stDateInput label{color:#94A3B8!important;font-size:0.75rem!important;}
[data-testid="stSidebar"] input{background:#0F172A!important;color:#E2E8F0!important;border-color:#334155!important;}
h1,h2,h3{font-family:'Inter',sans-serif!important;color:var(--txt)!important;}
p,span,div,label{font-family:'Inter',sans-serif!important;}
#MainMenu,footer,header{visibility:hidden;}
[data-testid="stToolbar"]{display:none;}
[data-testid="InputInstructions"]{display:none!important;}
kbd{display:none!important;}
/* keyboard_double / sidebar collapse button fix */
[role="tooltip"]{display:none!important;}
[data-baseweb="tooltip"]{display:none!important;}
.material-symbols-rounded{font-size:0!important;width:0!important;height:0!important;overflow:hidden!important;display:inline-block!important;}
[aria-label*="keyboard"]{display:none!important;}
div[class*="Tooltip"]{display:none!important;}
span[class*="instruction"]{display:none!important;}
[data-testid="InputInstructions"]{display:none!important;visibility:hidden!important;}
span[class*="material"]{font-size:0!important;color:transparent!important;width:0!important;}
[data-testid="stWidgetLabel"] span{font-size:0!important;}
[data-testid="stWidgetLabel"] span[data-testid="stWidgetLabelHelpInline"]{display:none!important;}
/* Forzar sidebar siempre visible y ocultar botón collapse */
[data-testid="stSidebar"]{display:block!important;visibility:visible!important;transform:none!important;left:0!important;min-width:244px!important;}
[data-testid="collapsedControl"]{display:none!important;}
[data-testid="stSidebarCollapseButton"]{display:none!important;}
button[aria-label="Close sidebar"]{display:none!important;}
button[aria-label="Open sidebar"]{display:none!important;}
section[data-testid="stSidebar"] > div:first-child > div:first-child > button{display:none!important;}
.kpi{background:var(--surf);border:1px solid var(--bord);border-radius:12px;padding:1.2rem 1.4rem;box-shadow:0 1px 3px rgba(0,0,0,0.06);}
.kpi-badge{display:inline-block;font-size:0.65rem;font-weight:700;letter-spacing:0.08em;text-transform:uppercase;padding:3px 10px;border-radius:20px;margin-bottom:0.7rem;}
.kpi-val{font-family:'IBM Plex Mono',monospace;font-size:2rem;font-weight:600;color:var(--txt);line-height:1;}
.kpi-mw{font-size:1rem;font-weight:400;color:var(--muted);}
.kpi-sub{font-size:0.75rem;color:var(--muted);margin-top:0.4rem;}
.kpi-delta{font-size:0.78rem;margin-top:0.5rem;font-weight:500;}
.sec{font-size:0.82rem;font-weight:800;letter-spacing:0.12em;text-transform:uppercase;color:#334155;border-bottom:2px solid var(--bord);padding-bottom:0.45rem;margin:1.8rem 0 1rem;}
.dot-status{display:inline-block;width:8px;height:8px;border-radius:50%;margin-right:6px;vertical-align:middle;}
.dot-g{background:#10B981;box-shadow:0 0 5px rgba(16,185,129,0.6);animation:blink 2s infinite;}
.dot-r{background:#EF4444;}
.dot-y{background:#F59E0B;animation:blink 2s infinite;}
@keyframes blink{0%,100%{opacity:1}50%{opacity:.3}}
@keyframes pulse-pend{0%,100%{box-shadow:0 0 0 0 rgba(217,119,6,0.7)}70%{box-shadow:0 0 0 6px rgba(217,119,6,0)}}
.badge-pend{display:inline-block;animation:pulse-pend 1.8s infinite;border-radius:4px;}
@keyframes pulse-sscc{0%,100%{box-shadow:0 0 0 0 rgba(100,116,139,0.5)}70%{box-shadow:0 0 0 5px rgba(100,116,139,0)}}
.sscc-latest{animation:pulse-sscc 2.2s infinite;}
.status-box{background:var(--surf2);border:1px solid var(--bord);border-radius:8px;padding:0.6rem 0.8rem;margin-top:0.5rem;font-size:0.72rem;}
.stTextInput>div>div>input,.stTextArea>div>div>textarea,.stNumberInput>div>div>input{
    background:var(--surf2)!important;border:1px solid var(--bord)!important;border-radius:8px!important;color:var(--txt)!important;}
.stButton>button{background:var(--accent)!important;color:#fff!important;border:none!important;border-radius:8px!important;font-family:'Inter',sans-serif!important;font-weight:600!important;}
.stButton>button:hover{opacity:.88!important;}
.stTabs [data-baseweb="tab-list"]{gap:4px;}
.stTabs [data-baseweb="tab"]{border-radius:8px 8px 0 0;font-weight:600;font-family:'Inter',sans-serif;}
/* ── Selectbox: valor mostrado ── */
[data-testid="stSelectbox"] div[data-baseweb="select"] > div{color:#0F172A!important;background:#FFFFFF!important;}
[data-testid="stSelectbox"] div[data-baseweb="select"] span{color:#0F172A!important;}
[data-testid="stSelectbox"] div[data-baseweb="select"] input{color:#0F172A!important;}
/* Sidebar selectbox — texto claro sobre fondo oscuro */
[data-testid="stSidebar"] [data-testid="stSelectbox"] div[data-baseweb="select"] > div{color:#E2E8F0!important;background:#1E293B!important;}
[data-testid="stSidebar"] [data-testid="stSelectbox"] div[data-baseweb="select"] span{color:#E2E8F0!important;}
/* ── Lista desplegable (portal baseweb, fuera del sidebar DOM) ── */
[data-baseweb="popover"] ul li div,
[data-baseweb="popover"] ul li span,
[data-baseweb="popover"] [role="option"] div,
[data-baseweb="popover"] [role="option"] span{color:#0F172A!important;}
[data-baseweb="popover"] ul{background:#FFFFFF!important;}
[data-baseweb="popover"] [role="option"]{background:#FFFFFF!important;}
[data-baseweb="popover"] [role="option"]:hover{background:#EFF6FF!important;}
[data-baseweb="popover"] [aria-selected="true"]{background:#DBEAFE!important;}
/* ── Select nativo (bitácora "Ver registros" filtro) ── */
select,select option{color:#0F172A!important;background:#FFFFFF!important;}
</style>
<script>
function hideKeyboardHints() {
    document.querySelectorAll('[data-testid="InputInstructions"]').forEach(function(el){el.style.display='none';});
    document.querySelectorAll('.material-symbols-rounded').forEach(function(el){el.style.fontSize='0';el.style.width='0';});
}
hideKeyboardHints();
setInterval(hideKeyboardHints, 500);

// Limpiar estado colapsado de sidebar en localStorage y sessionStorage
(function() {
    try {
        [localStorage, sessionStorage].forEach(function(store) {
            Object.keys(store).forEach(function(k) {
                if (/sidebar|Sidebar|collapsed/i.test(k)) store.removeItem(k);
            });
        });
    } catch(e) {}
})();

// MutationObserver: revertir cualquier intento de Streamlit de colapsar el sidebar
(function() {
    function forceSidebar() {
        var sb = document.querySelector('[data-testid="stSidebar"]');
        if (!sb) return;
        sb.style.setProperty('display','block','important');
        sb.style.setProperty('visibility','visible','important');
        sb.style.setProperty('transform','none','important');
        sb.style.setProperty('left','0','important');
        if (sb.getAttribute('data-collapsed') === 'true') {
            sb.setAttribute('data-collapsed','false');
        }
    }
    forceSidebar();
    var obs = new MutationObserver(forceSidebar);
    obs.observe(document.body, {subtree:true, attributes:true, childList:true});
})();

// Inyectar CSS en <head> para dropdown options (mayor especificidad que variables Streamlit)
(function injectDropdownCSS() {
    var style = document.createElement('style');
    style.innerHTML = [
        '[data-baseweb="popover"] [role="option"] { color: #0F172A !important; background-color: #ffffff !important; }',
        '[data-baseweb="popover"] [role="option"]:hover { background-color: #EFF6FF !important; color: #1E40AF !important; }',
        '[data-baseweb="popover"] [role="option"] * { color: #0F172A !important; }',
        '[data-baseweb="popover"] [aria-selected="true"] { background-color: #DBEAFE !important; }',
        '[data-baseweb="popover"] [aria-selected="true"] * { color: #1E40AF !important; }',
        '[data-baseweb="list"] { background-color: #ffffff !important; }',
        '[data-baseweb="list"] * { color: #0F172A !important; }',
        'ul[data-baseweb="menu"] li { color: #0F172A !important; background-color: #ffffff !important; }',
        'ul[data-baseweb="menu"] li * { color: #0F172A !important; }',
    ].join('\n');
    document.head.appendChild(style);
    // Re-aplicar cuando Streamlit monta nuevos portales
    var obs2 = new MutationObserver(function(muts) {
        muts.forEach(function(m) {
            m.addedNodes.forEach(function(node) {
                if (node.nodeType !== 1) return;
                var opts = node.querySelectorAll ? node.querySelectorAll('[role="option"]') : [];
                opts.forEach(function(el) {
                    el.style.setProperty('color','#0F172A','important');
                    el.style.setProperty('background-color','#ffffff','important');
                });
            });
        });
    });
    obs2.observe(document.body, {childList:true, subtree:true});
})();
</script>
""", unsafe_allow_html=True)



# ══════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════
# HELPERS COMUNES — gráficos matplotlib reutilizados por PDF y PPT
# ══════════════════════════════════════════════════════════════
_UNIT_COLORS = {
    "ANG1": {"real": "#7C3AED", "prog": "#C4B5FD"},
    "ANG2": {"real": "#2563EB", "prog": "#93C5FD"},
    "CCR1": {"real": "#CA8A04", "prog": "#FDE68A"},
    "CCR2": {"real": "#16A34A", "prog": "#86EFAC"},
}
_UNIT_NAMES = {"ANG1": "Angamos 1", "ANG2": "Angamos 2", "CCR1": "Cochrane 1", "CCR2": "Cochrane 2"}


def _fig_generacion(df_real, df_prog, unidad, figsize=(14, 4)):
    """Retorna figura matplotlib con real vs programada para una unidad."""
    import matplotlib.pyplot as _plt
    import matplotlib.dates as _mdates
    col = _UNIT_COLORS[unidad]
    fig, ax = _plt.subplots(figsize=figsize)
    df_u  = df_real[df_real["unidad"] == unidad].sort_values("fecha_hora") if not df_real.empty else pd.DataFrame()
    df_up = df_prog[df_prog["unidad"] == unidad].sort_values("fecha_hora") if not df_prog.empty else pd.DataFrame()
    if df_u.empty:
        ax.text(0.5, 0.5, "Sin datos", ha="center", va="center", transform=ax.transAxes, color="#94A3B8")
    else:
        ax.plot(df_u["fecha_hora"], df_u["gen_real_mw"],
                color=col["real"], linewidth=2.2, label="Real MW", zorder=3)
        if not df_up.empty:
            ax.plot(df_up["fecha_hora"], df_up["gen_programada_mw"],
                    color=col["prog"], linewidth=2.0, linestyle="--", label="Programada MW", zorder=2)
            df_m = pd.merge_asof(
                df_u[["fecha_hora", "gen_real_mw"]].sort_values("fecha_hora"),
                df_up[["fecha_hora", "gen_programada_mw"]].sort_values("fecha_hora"),
                on="fecha_hora", direction="nearest", tolerance=pd.Timedelta("1h")
            ).dropna()
            if not df_m.empty:
                ax.fill_between(df_m["fecha_hora"], df_m["gen_real_mw"], df_m["gen_programada_mw"],
                                alpha=0.15, color=col["real"])
    ax.set_ylabel("MW", fontsize=9)
    ax.set_ylim(bottom=0)
    ax.xaxis.set_major_formatter(_mdates.DateFormatter("%a %d/%m"))
    ax.xaxis.set_major_locator(_mdates.DayLocator())
    _plt.setp(ax.xaxis.get_majorticklabels(), fontsize=8, rotation=15)
    ax.tick_params(axis="y", labelsize=8)
    ax.set_facecolor("#FAFBFF")
    ax.grid(axis="y", color="#E2E8F0", linewidth=0.7)
    ax.grid(axis="x", color="#F1F5F9", linewidth=0.5)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(loc="lower right", fontsize=8, framealpha=0.85)
    fig.patch.set_facecolor("white")
    fig.tight_layout(pad=0.4)
    return fig


def _fig_cmg(df_cmg, figsize=(14, 3)):
    """Retorna figura matplotlib con CMG USD/MWh."""
    import matplotlib.pyplot as _plt
    import matplotlib.dates as _mdates
    fig, ax = _plt.subplots(figsize=figsize)
    if df_cmg.empty:
        ax.text(0.5, 0.5, "Sin datos CMG", ha="center", va="center", transform=ax.transAxes, color="#94A3B8")
    else:
        ax.plot(df_cmg["fecha_hora"], df_cmg["cmg_usd_mwh"], color="#0891B2", linewidth=1.8)
        ax.fill_between(df_cmg["fecha_hora"], df_cmg["cmg_usd_mwh"], alpha=0.12, color="#0891B2")
    ax.set_ylabel("USD/MWh", fontsize=9)
    ax.xaxis.set_major_formatter(_mdates.DateFormatter("%a %d/%m"))
    ax.xaxis.set_major_locator(_mdates.DayLocator())
    _plt.setp(ax.xaxis.get_majorticklabels(), fontsize=8, rotation=15)
    ax.tick_params(axis="y", labelsize=8)
    ax.set_facecolor("#FAFBFF")
    ax.grid(axis="y", color="#E2E8F0", linewidth=0.7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.patch.set_facecolor("white")
    fig.tight_layout(pad=0.4)
    return fig


def _fig_to_bytes(fig, dpi=150):
    import io as _io, matplotlib.pyplot as _plt
    buf = _io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    _plt.close(fig)
    buf.seek(0)
    return buf


def _kpis_unidad(df_real, df_prog, unidad):
    """Retorna dict con KPIs calculados para una unidad."""
    df_u  = df_real[df_real["unidad"] == unidad] if not df_real.empty else pd.DataFrame()
    df_up = df_prog[df_prog["unidad"] == unidad] if not df_prog.empty else pd.DataFrame()
    pmax  = {"ANG1": 277, "ANG2": 280, "CCR1": 276, "CCR2": 276}[unidad]
    if df_u.empty:
        return {"prom": "—", "max": "—", "min": "—", "fc": "—", "desv": "—"}
    prom = df_u["gen_real_mw"].mean()
    mx   = df_u["gen_real_mw"].max()
    mn   = df_u["gen_real_mw"].min()
    fc   = prom / pmax * 100
    desv = "—"
    if not df_up.empty:
        df_m = pd.merge_asof(
            df_u[["fecha_hora", "gen_real_mw"]].sort_values("fecha_hora"),
            df_up[["fecha_hora", "gen_programada_mw"]].sort_values("fecha_hora"),
            on="fecha_hora", direction="nearest", tolerance=pd.Timedelta("1h")
        ).dropna()
        if not df_m.empty:
            desv = f"{(df_m['gen_real_mw'] - df_m['gen_programada_mw']).mean():+.1f}"
    return {
        "prom": f"{prom:.1f}", "max": f"{mx:.1f}", "min": f"{mn:.1f}",
        "fc":   f"{fc:.1f}%",  "desv": desv,
    }


# ══════════════════════════════════════════════════════════════
# GENERADOR PDF — Reporte Operacional Completo
# ══════════════════════════════════════════════════════════════
def generar_pdf(df_real, df_prog, df_cmg, start_str, end_str,
                df_sscc=None, df_lim=None):
    from reportlab.lib.pagesizes import A4, landscape as RL_landscape
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor, white
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Image as RLImage, HRFlowable, PageBreak, Table, TableStyle)
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    import io as _io

    if df_sscc is None: df_sscc = pd.DataFrame()
    if df_lim  is None: df_lim  = pd.DataFrame()

    buf = _io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=1.8*cm, rightMargin=1.8*cm,
                            topMargin=1.5*cm, bottomMargin=1.5*cm)

    # Paleta
    C_DARK  = HexColor("#0F172A")
    C_BLUE  = HexColor("#2563EB")
    C_GRAY  = HexColor("#475569")
    C_LGRAY = HexColor("#94A3B8")
    C_BG    = HexColor("#F8FAFC")
    C_LINE  = HexColor("#E2E8F0")
    C_GREEN = HexColor("#16A34A")
    C_AMBER = HexColor("#D97706")

    # Estilos
    def _sty(name, font="Helvetica", size=10, color=None, after=4, align=TA_LEFT, bold=False):
        return ParagraphStyle(name, fontName=f"Helvetica-Bold" if bold else font,
                              fontSize=size, textColor=color or C_DARK,
                              spaceAfter=after, alignment=align)

    sTitle  = _sty("tt", size=20, bold=True, after=4)
    sSub    = _sty("ss", size=11, color=C_GRAY, after=3)
    sH2     = _sty("h2", size=13, bold=True, after=6, color=C_BLUE)
    sH3     = _sty("h3", size=11, bold=True, after=4)
    sBody   = _sty("bd", size=9,  color=C_GRAY, after=3)
    sSmall  = _sty("sm", size=8,  color=C_LGRAY, after=2)
    sCenter = _sty("ct", size=8,  color=C_LGRAY, after=2, align=TA_CENTER)
    sBit    = _sty("bt", size=9,  color=C_DARK, after=2)

    try:
        dt_s = datetime.strptime(start_str, "%Y-%m-%d")
        semana_num = dt_s.isocalendar()[1]
        fmt_s = dt_s.strftime("%d/%m/%Y")
        fmt_e = datetime.strptime(end_str, "%Y-%m-%d").strftime("%d/%m/%Y")
    except:
        semana_num, fmt_s, fmt_e = "—", start_str, end_str

    story = []

    # ── PORTADA ───────────────────────────────────────────────
    story += [
        Spacer(1, 2.5*cm),
        Paragraph("Complejo Térmico Mejillones", sTitle),
        Paragraph("AES Andes · Monitoreo Operacional", sSub),
        Spacer(1, 0.8*cm),
        HRFlowable(width="100%", thickness=1.5, color=C_BLUE),
        Spacer(1, 0.6*cm),
        Paragraph(f"Reporte Operacional — Semana {semana_num}", _sty("rw", size=15, bold=True, after=4)),
        Paragraph(f"Período: {fmt_s} al {fmt_e}", _sty("rp", size=12, color=C_GRAY, after=8)),
        Spacer(1, 0.6*cm),
    ]

    # Resumen ejecutivo portada — KPIs del complejo
    resumen_rows = [["Unidad", "Prom Real (MW)", "Máx (MW)", "Mín (MW)", "Factor Carga", "Desv. vs Prog."]]
    for u in ["ANG1", "ANG2", "CCR1", "CCR2"]:
        k = _kpis_unidad(df_real, df_prog, u)
        resumen_rows.append([_UNIT_NAMES[u], k["prom"], k["max"], k["min"], k["fc"], k["desv"]])

    tbl_res = Table(resumen_rows, colWidths=[4*cm, 3*cm, 2.5*cm, 2.5*cm, 3*cm, 3*cm])
    tbl_res.setStyle(TableStyle([
        ("BACKGROUND",  (0,0), (-1,0), C_BLUE),
        ("TEXTCOLOR",   (0,0), (-1,0), white),
        ("FONTNAME",    (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",    (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [HexColor("#F8FAFC"), white]),
        ("GRID",        (0,0), (-1,-1), 0.5, C_LINE),
        ("ALIGN",       (1,0), (-1,-1), "CENTER"),
        ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",  (0,0), (-1,-1), 5),
        ("BOTTOMPADDING",(0,0),(-1,-1), 5),
    ]))
    story += [tbl_res, Spacer(1, 0.6*cm)]

    # CMG resumen portada
    if not df_cmg.empty:
        story.append(Paragraph(
            f"CMG promedio período — Crucero 220kV: <b>{df_cmg['cmg_usd_mwh'].mean():.1f} USD/MWh</b>"
            f"  ·  Mín: {df_cmg['cmg_usd_mwh'].min():.1f}  ·  Máx: {df_cmg['cmg_usd_mwh'].max():.1f}",
            sBody
        ))

    # Limitaciones activas en portada
    if not df_lim.empty:
        n_act = int((df_lim["status"] == "pendiente").sum())
        n_sscc_lim = int(df_lim["afecta_sscc"].fillna(False).sum())
        story.append(Paragraph(
            f"Limitaciones de transmisión — <b>{n_act} activas</b> en el período"
            f"  ·  {n_sscc_lim} afectan SSCC",
            sBody
        ))

    story += [
        Spacer(1, 1.5*cm),
        HRFlowable(width="100%", thickness=0.5, color=C_LINE),
        Paragraph(f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}  ·  Fuente: API CEN SIPUB / OPS / CMG S3",
                  sCenter),
        PageBreak(),
    ]

    # ── PÁGINA CMG ────────────────────────────────────────────
    story += [Paragraph("Costo Marginal (CMG)", sH2), Spacer(1, 0.2*cm)]
    fig_cmg = _fig_cmg(df_cmg, figsize=(14, 3.5))
    story.append(RLImage(_fig_to_bytes(fig_cmg), width=17*cm, height=6*cm))
    if not df_cmg.empty:
        story.append(Paragraph(
            f"Nodo Crucero 220kV (preliminar)  ·  Prom: {df_cmg['cmg_usd_mwh'].mean():.1f}  "
            f"Mín: {df_cmg['cmg_usd_mwh'].min():.1f}  Máx: {df_cmg['cmg_usd_mwh'].max():.1f} USD/MWh",
            sSmall
        ))
    story += [Spacer(1, 0.5*cm), PageBreak()]

    # ── PÁGINAS POR UNIDAD ────────────────────────────────────
    for u in ["ANG1", "ANG2", "CCR1", "CCR2"]:
        k = _kpis_unidad(df_real, df_prog, u)
        story.append(Paragraph(_UNIT_NAMES[u], sH2))

        # KPIs en tabla horizontal
        kpi_data = [
            ["Prom Real", "Máx", "Mín", "Factor Carga", "Desv. vs Prog."],
            [f"{k['prom']} MW", f"{k['max']} MW", f"{k['min']} MW", k["fc"], f"{k['desv']} MW"],
        ]
        tbl_kpi = Table(kpi_data, colWidths=[3.4*cm]*5)
        tbl_kpi.setStyle(TableStyle([
            ("BACKGROUND",  (0,0), (-1,0), HexColor("#EFF6FF")),
            ("FONTNAME",    (0,0), (-1,0), "Helvetica-Bold"),
            ("FONTNAME",    (0,1), (-1,1), "Helvetica-Bold"),
            ("FONTSIZE",    (0,0), (-1,-1), 9),
            ("ALIGN",       (0,0), (-1,-1), "CENTER"),
            ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
            ("GRID",        (0,0), (-1,-1), 0.5, C_LINE),
            ("TOPPADDING",  (0,0), (-1,-1), 5),
            ("BOTTOMPADDING",(0,0),(-1,-1), 5),
            ("TEXTCOLOR",   (0,1), (-1,1), C_BLUE),
        ]))
        story += [tbl_kpi, Spacer(1, 0.3*cm)]

        # Gráfico gen real vs programada
        fig_g = _fig_generacion(df_real, df_prog, u, figsize=(14, 4))
        story.append(RLImage(_fig_to_bytes(fig_g), width=17*cm, height=7*cm))
        story.append(Spacer(1, 0.3*cm))

        # SSCC de la unidad
        if not df_sscc.empty:
            df_su = df_sscc[df_sscc["unidad"] == u]
            if not df_su.empty:
                story.append(Paragraph(f"SSCC — {len(df_su)} instrucciones en el período", sH3))
                sscc_rows = [["Fecha", "Inicio", "Fin", "Instrucción", "Motivo"]]
                for _, row in df_su.head(10).iterrows():
                    sscc_rows.append([
                        str(row.get("fecha", ""))[:10],
                        str(row.get("inicio_periodo", ""))[:5],
                        str(row.get("fin_periodo", ""))[:5],
                        str(row.get("instruccion_sscc", "")),
                        str(row.get("motivo", "") or "")[:60],
                    ])
                tbl_sscc = Table(sscc_rows, colWidths=[2.2*cm, 1.8*cm, 1.8*cm, 2.5*cm, 8.7*cm])
                tbl_sscc.setStyle(TableStyle([
                    ("BACKGROUND",  (0,0), (-1,0), HexColor("#F0FDF4")),
                    ("FONTNAME",    (0,0), (-1,0), "Helvetica-Bold"),
                    ("FONTSIZE",    (0,0), (-1,-1), 8),
                    ("GRID",        (0,0), (-1,-1), 0.4, C_LINE),
                    ("ROWBACKGROUNDS", (0,1), (-1,-1), [white, HexColor("#F8FAFC")]),
                    ("TOPPADDING",  (0,0), (-1,-1), 3),
                    ("BOTTOMPADDING",(0,0),(-1,-1), 3),
                ]))
                story += [tbl_sscc, Spacer(1, 0.2*cm)]

        # Limitaciones de la unidad
        if not df_lim.empty and "_unidad" in df_lim.columns:
            df_lu = df_lim[df_lim["_unidad"] == u]
            if not df_lu.empty:
                story.append(Paragraph(f"Limitaciones — {len(df_lu)} en el período", sH3))
                lim_rows = [["Correlativo", "Status", "Apertura", "Retorno est.", "Potencia", "Afecta SSCC"]]
                for _, row in df_lu.head(8).iterrows():
                    corr = str(int(float(row["correlativo"]))) if pd.notna(row.get("correlativo")) else "—"
                    pot  = f"{int(float(row['potencia']))} MW" if pd.notna(row.get("potencia")) and float(row["potencia"]) > 0 else "—"
                    lim_rows.append([
                        corr,
                        str(row.get("status", "")),
                        str(row.get("fecha_perturbacion", ""))[:16],
                        str(row.get("fecha_retorno_estimada", "") or "—")[:10],
                        pot,
                        "Sí" if row.get("afecta_sscc") else "No",
                    ])
                tbl_lim = Table(lim_rows, colWidths=[2.5*cm, 2.2*cm, 3.5*cm, 2.5*cm, 2.5*cm, 2.8*cm])
                tbl_lim.setStyle(TableStyle([
                    ("BACKGROUND",  (0,0), (-1,0), HexColor("#FFFBEB")),
                    ("FONTNAME",    (0,0), (-1,0), "Helvetica-Bold"),
                    ("FONTSIZE",    (0,0), (-1,-1), 8),
                    ("GRID",        (0,0), (-1,-1), 0.4, C_LINE),
                    ("ROWBACKGROUNDS", (0,1), (-1,-1), [white, HexColor("#F8FAFC")]),
                    ("TOPPADDING",  (0,0), (-1,-1), 3),
                    ("BOTTOMPADDING",(0,0),(-1,-1), 3),
                ]))
                story += [tbl_lim, Spacer(1, 0.2*cm)]

        # Bitácora
        try:
            novedades = qry(
                "SELECT fecha, hora, comentario FROM bitacora "
                "WHERE unidad=%s AND fecha BETWEEN %s AND %s ORDER BY fecha, hora",
                (u, start_str, end_str)
            )
        except: novedades = pd.DataFrame()

        if not novedades.empty:
            story.append(Paragraph("Bitácora", sH3))
            for _, row in novedades.iterrows():
                try: fd = datetime.strptime(str(row["fecha"]), "%Y-%m-%d").strftime("%d/%m/%Y")
                except: fd = str(row["fecha"])
                story.append(Paragraph(
                    f"<b>{fd} {str(row['hora'])[:5]}</b> — {str(row['comentario'])}", sBit
                ))

        story += [
            Spacer(1, 0.4*cm),
            HRFlowable(width="100%", thickness=0.4, color=C_LINE),
            Paragraph(f"Reporte generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}", sCenter),
            PageBreak(),
        ]

    # Quitar último PageBreak
    if story and isinstance(story[-1], PageBreak):
        story.pop()

    doc.build(story)
    buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════
# GENERADOR PPT — Presentación Operacional
# ══════════════════════════════════════════════════════════════
def generar_ppt(df_real, df_prog, df_cmg, start_str, end_str,
                df_sscc=None, df_lim=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io as _io

    if df_sscc is None: df_sscc = pd.DataFrame()
    if df_lim  is None: df_lim  = pd.DataFrame()

    # Colores
    RGB_DARK  = RGBColor(0x0F, 0x17, 0x2A)
    RGB_BLUE  = RGBColor(0x25, 0x63, 0xEB)
    RGB_GRAY  = RGBColor(0x47, 0x55, 0x69)
    RGB_LGRAY = RGBColor(0x94, 0xA3, 0xB8)
    RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    RGB_BG    = RGBColor(0xF8, 0xFA, 0xFC)
    RGB_AMBER = RGBColor(0xD9, 0x77, 0x06)
    RGB_GREEN = RGBColor(0x16, 0xA3, 0x4A)

    try:
        dt_s = datetime.strptime(start_str, "%Y-%m-%d")
        semana_num = dt_s.isocalendar()[1]
        fmt_s = dt_s.strftime("%d/%m/%Y")
        fmt_e = datetime.strptime(end_str, "%Y-%m-%d").strftime("%d/%m/%Y")
    except:
        semana_num, fmt_s, fmt_e = "—", start_str, end_str

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # completamente en blanco

    def _add_slide():
        return prs.slides.add_slide(blank)

    def _bg(slide, color=RGB_WHITE):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _txb(slide, text, l, t, w, h, size=18, bold=False, color=RGB_DARK, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def _rect(slide, l, t, w, h, color=RGB_BLUE):
        from pptx.util import Inches as _I
        shp = slide.shapes.add_shape(1, _I(l), _I(t), _I(w), _I(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def _img(slide, fig, l, t, w, h):
        buf = _fig_to_bytes(fig, dpi=150)
        slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))

    # ── SLIDE 1: PORTADA ──────────────────────────────────────
    sl = _add_slide()
    _bg(sl, RGB_DARK)
    _rect(sl, 0, 0, 13.33, 0.12, RGB_BLUE)
    _rect(sl, 0, 7.38, 13.33, 0.12, RGB_BLUE)
    _txb(sl, "Complejo Térmico Mejillones", 0.6, 1.8, 12, 1.0,
         size=36, bold=True, color=RGB_WHITE)
    _txb(sl, "AES Andes · Monitoreo Operacional", 0.6, 2.8, 12, 0.6,
         size=18, color=RGB_LGRAY)
    _txb(sl, f"Reporte Semana {semana_num}  ·  {fmt_s} — {fmt_e}", 0.6, 3.5, 12, 0.5,
         size=16, color=RGB_BLUE)
    _txb(sl, f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}", 0.6, 6.9, 12, 0.4,
         size=10, color=RGB_LGRAY)

    # ── SLIDE 2: RESUMEN EJECUTIVO ─────────────────────────────
    sl = _add_slide()
    _bg(sl, RGB_BG)
    _rect(sl, 0, 0, 13.33, 0.55, RGB_BLUE)
    _txb(sl, "Resumen Ejecutivo", 0.3, 0.08, 10, 0.4, size=18, bold=True, color=RGB_WHITE)

    col_x = [0.3, 3.5, 6.3, 9.1]
    unidades = ["ANG1", "ANG2", "CCR1", "CCR2"]
    col_colors = [RGBColor(0x7C,0x3A,0xED), RGBColor(0x25,0x63,0xEB),
                  RGBColor(0xCA,0x8A,0x04), RGBColor(0x16,0xA3,0x4A)]
    for i, (u, cx, cc) in enumerate(zip(unidades, col_x, col_colors)):
        k = _kpis_unidad(df_real, df_prog, u)
        _rect(sl, cx, 0.7, 2.9, 0.35, cc)
        _txb(sl, _UNIT_NAMES[u], cx+0.1, 0.73, 2.7, 0.3, size=13, bold=True, color=RGB_WHITE)
        items = [
            f"Prom real:  {k['prom']} MW",
            f"Máx:  {k['max']} MW",
            f"Mín:  {k['min']} MW",
            f"Factor carga:  {k['fc']}",
            f"Desv. prog.:  {k['desv']} MW",
        ]
        for j, txt in enumerate(items):
            _txb(sl, txt, cx+0.1, 1.15 + j*0.48, 2.8, 0.45, size=11, color=RGB_DARK)

    # CMG + Limitaciones resumen abajo
    if not df_cmg.empty:
        cmg_txt = (f"CMG Crucero 220kV  ·  Prom: {df_cmg['cmg_usd_mwh'].mean():.1f}  "
                   f"Mín: {df_cmg['cmg_usd_mwh'].min():.1f}  Máx: {df_cmg['cmg_usd_mwh'].max():.1f} USD/MWh")
        _txb(sl, cmg_txt, 0.3, 6.7, 8, 0.4, size=10, color=RGB_GRAY)
    if not df_lim.empty:
        n_act = int((df_lim["status"] == "pendiente").sum())
        _txb(sl, f"Limitaciones activas: {n_act}", 9.0, 6.7, 4, 0.4, size=10,
             bold=(n_act > 0), color=RGB_AMBER if n_act > 0 else RGB_GRAY)

    # ── SLIDE 3: CMG ──────────────────────────────────────────
    sl = _add_slide()
    _bg(sl, RGB_BG)
    _rect(sl, 0, 0, 13.33, 0.55, RGBColor(0x08, 0x91, 0xB2))
    _txb(sl, "Costo Marginal — CMG", 0.3, 0.08, 10, 0.4, size=18, bold=True, color=RGB_WHITE)
    fig_cmg = _fig_cmg(df_cmg, figsize=(13, 4.5))
    _img(sl, fig_cmg, 0.3, 0.7, 12.7, 5.5)
    if not df_cmg.empty:
        _txb(sl, f"Nodo Crucero 220kV (preliminar)  ·  Prom: {df_cmg['cmg_usd_mwh'].mean():.1f}  "
             f"Mín: {df_cmg['cmg_usd_mwh'].min():.1f}  Máx: {df_cmg['cmg_usd_mwh'].max():.1f} USD/MWh",
             0.3, 6.4, 12, 0.5, size=10, color=RGB_GRAY)

    # ── SLIDES POR UNIDAD ─────────────────────────────────────
    for u, cc in zip(unidades, col_colors):
        k = _kpis_unidad(df_real, df_prog, u)

        # Slide generación
        sl = _add_slide()
        _bg(sl, RGB_BG)
        _rect(sl, 0, 0, 13.33, 0.55, cc)
        _txb(sl, f"{_UNIT_NAMES[u]} — Generación", 0.3, 0.08, 10, 0.4,
             size=18, bold=True, color=RGB_WHITE)

        # KPIs en banda
        kpi_labels = ["Prom Real", "Máximo", "Mínimo", "Factor Carga", "Desv. Prog."]
        kpi_vals   = [f"{k['prom']} MW", f"{k['max']} MW", f"{k['min']} MW", k["fc"], f"{k['desv']} MW"]
        for i, (lbl, val) in enumerate(zip(kpi_labels, kpi_vals)):
            bx = 0.3 + i * 2.6
            _rect(sl, bx, 0.65, 2.4, 0.28, RGBColor(0xEF,0xF6,0xFF))
            _txb(sl, lbl, bx+0.05, 0.67, 2.3, 0.22, size=8, color=RGB_GRAY, align=PP_ALIGN.CENTER)
            _txb(sl, val,  bx+0.05, 0.88, 2.3, 0.28, size=11, bold=True, color=cc, align=PP_ALIGN.CENTER)

        fig_g = _fig_generacion(df_real, df_prog, u, figsize=(13, 4.2))
        _img(sl, fig_g, 0.3, 1.22, 12.7, 5.2)

        # Slide SSCC + Limitaciones + Bitácora
        sl = _add_slide()
        _bg(sl, RGB_BG)
        _rect(sl, 0, 0, 13.33, 0.55, cc)
        _txb(sl, f"{_UNIT_NAMES[u]} — SSCC / Limitaciones / Bitácora",
             0.3, 0.08, 12, 0.4, size=18, bold=True, color=RGB_WHITE)

        cur_y = 0.7

        # SSCC
        if not df_sscc.empty:
            df_su = df_sscc[df_sscc["unidad"] == u].head(6)
            if not df_su.empty:
                _txb(sl, f"SSCC ({len(df_sscc[df_sscc['unidad']==u])} instrucciones)", 0.3, cur_y, 6, 0.3,
                     size=11, bold=True, color=RGB_GREEN)
                cur_y += 0.32
                for _, row in df_su.iterrows():
                    txt = (f"{str(row.get('fecha',''))[:10]}  {str(row.get('inicio_periodo',''))[:5]}–"
                           f"{str(row.get('fin_periodo',''))[:5]}  [{row.get('instruccion_sscc','')}]  "
                           f"{str(row.get('motivo','') or '')[:55]}")
                    _txb(sl, txt, 0.4, cur_y, 12.5, 0.3, size=9, color=RGB_DARK)
                    cur_y += 0.3

        cur_y += 0.1

        # Limitaciones
        if not df_lim.empty and "_unidad" in df_lim.columns:
            df_lu = df_lim[df_lim["_unidad"] == u].head(5)
            if not df_lu.empty:
                _txb(sl, f"Limitaciones ({len(df_lim[df_lim['_unidad']==u])})", 0.3, cur_y, 6, 0.3,
                     size=11, bold=True, color=RGB_AMBER)
                cur_y += 0.32
                for _, row in df_lu.iterrows():
                    corr = str(int(float(row["correlativo"]))) if pd.notna(row.get("correlativo")) else "—"
                    pot  = f"{int(float(row['potencia']))} MW" if pd.notna(row.get("potencia")) and float(row["potencia"]) > 0 else ""
                    txt  = (f"N.{corr}  [{row.get('status','')}]  "
                            f"{str(row.get('fecha_perturbacion',''))[:16]}  {pot}  "
                            f"{'[Afecta SSCC]' if row.get('afecta_sscc') else ''}")
                    _txb(sl, txt, 0.4, cur_y, 12.5, 0.3, size=9, color=RGB_DARK)
                    cur_y += 0.3

        cur_y += 0.1

        # Bitácora
        try:
            novedades = qry(
                "SELECT fecha, hora, comentario FROM bitacora "
                "WHERE unidad=%s AND fecha BETWEEN %s AND %s ORDER BY fecha, hora",
                (u, start_str, end_str)
            )
        except: novedades = pd.DataFrame()

        if not novedades.empty and cur_y < 6.8:
            _txb(sl, "Bitácora", 0.3, cur_y, 6, 0.3, size=11, bold=True, color=RGB_DARK)
            cur_y += 0.32
            for _, row in novedades.head(4).iterrows():
                try: fd = datetime.strptime(str(row["fecha"]), "%Y-%m-%d").strftime("%d/%m/%Y")
                except: fd = str(row["fecha"])
                _txb(sl, f"{fd} {str(row['hora'])[:5]} — {str(row['comentario'])[:90]}",
                     0.4, cur_y, 12.5, 0.32, size=9, color=RGB_DARK)
                cur_y += 0.32

    # ── SLIDE FINAL: LIMITACIONES RESUMEN ─────────────────────
    if not df_lim.empty:
        sl = _add_slide()
        _bg(sl, RGB_BG)
        _rect(sl, 0, 0, 13.33, 0.55, RGB_AMBER)
        _txb(sl, "Limitaciones de Transmisión — Resumen", 0.3, 0.08, 12, 0.4,
             size=18, bold=True, color=RGB_WHITE)
        cur_y = 0.7
        for _, row in df_lim.head(12).iterrows():
            corr   = str(int(float(row["correlativo"]))) if pd.notna(row.get("correlativo")) else "—"
            unidad = row.get("_unidad", "") if "_unidad" in df_lim.columns else ""
            pot    = f"{int(float(row['potencia']))} MW" if pd.notna(row.get("potencia")) and float(row["potencia"]) > 0 else "—"
            txt = (f"N.{corr}  {unidad}  [{row.get('status','')}]  "
                   f"Apertura: {str(row.get('fecha_perturbacion',''))[:16]}  "
                   f"Potencia: {pot}  "
                   f"{'[Afecta SSCC]' if row.get('afecta_sscc') else ''}")
            _txb(sl, txt, 0.3, cur_y, 12.7, 0.35, size=9, color=RGB_DARK)
            cur_y += 0.38
            if cur_y > 7.0: break

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── DB ────────────────────────────────────────────────────────
@st.cache_resource
def get_conn():
    url = st.secrets.get("DATABASE_URL","")
    if not url: st.error("DATABASE_URL no configurada."); st.stop()
    return psycopg2.connect(url)

def test_conn():
    try:
        conn = get_conn()
        with conn.cursor() as c: c.execute("SELECT 1")
        return True, None
    except Exception:
        get_conn.clear()
        try:
            conn = get_conn()
            with conn.cursor() as c: c.execute("SELECT 1")
            return True, None
        except Exception as e:
            return False, str(e)

def qry(sql, params=None):
    try:
        conn = get_conn()
        return pd.read_sql(sql, conn, params=params)
    except Exception:
        get_conn.clear()
        try: return pd.read_sql(sql, get_conn(), params=params)
        except Exception as ex: st.error(f"Error DB: {ex}"); return pd.DataFrame()

def exe(sql, params=None):
    try:
        conn = get_conn()
        with conn.cursor() as c: c.execute(sql, params)
        conn.commit(); return True
    except Exception as e: st.error(f"Error DB: {e}"); return False

def exe_many(sql, params_list):
    try:
        conn = get_conn()
        with conn.cursor() as c: c.executemany(sql, params_list)
        conn.commit(); return True
    except Exception as e: st.error(f"Error DB: {e}"); return False


# ── Datos ─────────────────────────────────────────────────────
@st.cache_data(ttl=300)
def load_real(s,e):
    df = qry("SELECT unidad,gen_real_mw,fecha_hora,hora,potencia_maxima FROM generacion_real WHERE fecha_hora::date BETWEEN %s AND %s ORDER BY unidad,fecha_hora",(s,e))
    if not df.empty: df["fecha_hora"]=pd.to_datetime(df["fecha_hora"])
    return df

@st.cache_data(ttl=300)
def load_prog(s,e):
    df = qry("""
        SELECT DISTINCT ON (unidad, fecha_hora)
            unidad, gen_programada_mw, fecha_hora, hora, fuente
        FROM generacion_programada
        WHERE fecha_hora::date BETWEEN %s AND %s
        ORDER BY unidad, fecha_hora,
            CASE fuente WHEN 'CEN_PCP' THEN 0 ELSE 1 END
    """, (s,e))
    if not df.empty: df["fecha_hora"]=pd.to_datetime(df["fecha_hora"])
    return df

@st.cache_data(ttl=300)
def load_cmg(s, e, nodo="CRUCERO_______220"):
    df = qry(
        "SELECT fecha_hora,hora,cmg_usd_mwh FROM costo_marginal "
        "WHERE barra_transf=%s AND fecha_hora::date BETWEEN %s AND %s "
        "ORDER BY fecha_hora",
        (nodo, s, e),
    )
    if not df.empty: df["fecha_hora"]=pd.to_datetime(df["fecha_hora"])
    return df

@st.cache_data(ttl=300)
def load_sscc(s, e):
    return qry(
        "SELECT unidad, instruccion_sscc, fecha, inicio_periodo, fin_periodo, "
        "disponibilidad, motivo, comentario, estado_sabana "
        "FROM sscc_instrucciones "
        "WHERE fecha BETWEEN %s AND %s "
        "ORDER BY fecha DESC, unidad, inicio_periodo",
        (s, e),
    )

@st.cache_data(ttl=300)
def load_limitaciones(s, e):
    # Muestra limitaciones cuya perturbación cae en el período
    # o que siguen activas (sin retorno efectivo) durante el período
    return qry(
        "SELECT id, correlativo, instalacion_nombre, status, fecha_perturbacion, "
        "fecha_retorno_estimada, fecha_efectiva_retorno, potencia, "
        "unidad_medida_potencia, afecta_sscc, observacion, id_central, id_unidad "
        "FROM limitaciones_transmision "
        "WHERE fecha_perturbacion::date BETWEEN %s AND %s "
        "   OR (fecha_perturbacion::date < %s AND "
        "       (fecha_efectiva_retorno IS NULL OR fecha_efectiva_retorno::date >= %s)) "
        "ORDER BY fecha_perturbacion DESC",
        (s, e, s, s),
    )

@st.cache_data(ttl=300)
def load_solicitudes(s, e):
    return qry(
        "SELECT id, correlativo, empresa_nombre, instalacion_nombre, status, "
        "tipo_solicitud, type, tipo_programacion, descripcion_nivel_riesgo, "
        "fecha_inicio, fecha_fin, modified "
        "FROM solicitudes_trabajo "
        "WHERE fecha_inicio::date <= %s AND fecha_fin::date >= %s "
        "   OR partition_date::date BETWEEN %s AND %s "
        "ORDER BY fecha_inicio DESC",
        (e, s, s, e),
    )

@st.cache_data(ttl=60)
def load_bit(s,e,u=None):
    if u and u!="Todas":
        return qry("SELECT id,unidad,autor,comentario,fecha,hora FROM bitacora WHERE fecha BETWEEN %s AND %s AND unidad=%s ORDER BY fecha DESC,hora DESC",(s,e,u))
    return qry("SELECT id,unidad,autor,comentario,fecha,hora FROM bitacora WHERE fecha BETWEEN %s AND %s ORDER BY fecha DESC,hora DESC",(s,e))


# ── Sidebar ───────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### Complejo Térmico Mejillones")
    st.markdown('<p style="font-size:0.75rem;color:#64748B;margin-bottom:0.3rem">Complejo Térmico · Monitoreo Operacional</p>', unsafe_allow_html=True)

    # Estado de conexión y fuentes
    db_ok, db_err = test_conn()
    dot_db = "dot-g" if db_ok else "dot-r"
    txt_db = "Conectado · Supabase / PostgreSQL" if db_ok else "Error de conexión DB"

    # Última adquisición por fuente
    _ult_r   = qry("SELECT MAX(fecha_hora) AS t FROM generacion_real")
    _ult_p   = qry("SELECT MAX(fecha_hora) AS t FROM generacion_programada WHERE fuente='CEN_PCP'")
    _ult_cmg = qry("SELECT MAX(fecha_hora) AS t FROM costo_marginal")
    _ult_s   = qry("SELECT MAX(fecha_accion) AS t FROM sscc_instrucciones")
    _ult_lim = qry("SELECT MAX(modified) AS t FROM limitaciones_transmision")
    _ult_sol = qry("SELECT MAX(modified) AS t FROM solicitudes_trabajo")

    def _fmt(df, fmt="%d/%m %H:%M"):
        v = df.iloc[0]["t"] if not df.empty else None
        if v is None or pd.isna(v): return "—"
        if hasattr(v, "strftime"): return v.strftime(fmt)
        return str(v)[:16]

    str_r   = _fmt(_ult_r)
    str_p   = _fmt(_ult_p)
    str_cmg = _fmt(_ult_cmg)
    str_s   = _fmt(_ult_s, "%d/%m/%Y")
    str_lim = _fmt(_ult_lim, "%d/%m/%Y")
    str_sol = _fmt(_ult_sol, "%d/%m/%Y")

    st.markdown(f"""
    <div class="status-box">
        <div style="margin-bottom:6px">
            <span class="dot-status {dot_db}"></span>
            <span style="font-size:0.72rem;font-weight:600">{txt_db}</span>
        </div>
        <div style="font-size:0.68rem;line-height:2">
            <span style="font-size:0.62rem;font-weight:700;letter-spacing:0.06em;color:#94A3B8;text-transform:uppercase">API CEN SIPUB / OPS</span><br>
            <span class="dot-status dot-g"></span>Gen. real → <b>{str_r}</b><br>
            <span class="dot-status dot-g"></span>Gen. programada → <b>{str_p}</b><br>
            <span class="dot-status dot-g"></span>CMG S3 → <b>{str_cmg}</b><br>
            <span class="dot-status dot-g"></span>SSCC → <b>{str_s}</b><br>
            <span class="dot-status dot-g"></span>Limitaciones → <b>{str_lim}</b><br>
            <span class="dot-status dot-g"></span>Solicitudes → <b>{str_sol}</b><br>
            <span class="dot-status dot-g"></span>Adquisición GitHub Actions · /hora
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")
    hoy   = date.today()
    inicio_semana = hoy - timedelta(days=7)
    fi = st.date_input("Desde", value=inicio_semana, max_value=hoy)
    ff = st.date_input("Hasta", value=hoy,   max_value=hoy)
    if fi > ff: st.error("Fecha inicio > fin"); st.stop()
    st.markdown("---")
    mostrar_prog = st.checkbox("Mostrar programada", value=True)
    mostrar_cmg  = st.checkbox("Mostrar CMG en gráficos", value=True)
    nodo_cmg = "CRUCERO_______220"
    if mostrar_cmg:
        nodo_sel = st.radio(
            "Nodo CMG",
            list(NOMBRES_NODO.keys()),
            format_func=lambda x: NOMBRES_NODO[x],
            key="nodo_cmg_sel",
        )
        nodo_cmg = nodo_sel
    st.markdown("---")
    if st.button("Actualizar datos"):
        st.cache_data.clear(); st.rerun()

    st.markdown("---")
    st.markdown('<p style="font-size:0.72rem;font-weight:600;color:#475569">EXPORTAR REPORTE</p>', unsafe_allow_html=True)

    def _datos_reporte():
        s_r = fi.strftime("%Y-%m-%d")
        e_r = ff.strftime("%Y-%m-%d")
        df_lim_r = load_limitaciones(s_r, e_r)
        if not df_lim_r.empty:
            df_lim_r["_unidad"] = df_lim_r["id_unidad"].apply(
                lambda x: {1965:"ANG1",1966:"ANG2",1967:"CCR1",1968:"CCR2"}.get(int(float(x)),"") if pd.notna(x) else ""
            )
        return (load_real(s_r, e_r), load_prog(s_r, e_r),
                load_cmg(s_r, e_r), load_sscc(s_r, e_r),
                df_lim_r, s_r, e_r)

    if st.button("Generar PDF"):
        with st.spinner("Generando PDF..."):
            try:
                dr, dp, dc, ds, dl, s_r, e_r = _datos_reporte()
                pdf_bytes = generar_pdf(dr, dp, dc, s_r, e_r, df_sscc=ds, df_lim=dl)
                st.download_button(
                    "Descargar PDF",
                    data=pdf_bytes,
                    file_name=f"CTM-Reporte_{s_r}_{e_r}.pdf",
                    mime="application/pdf",
                )
            except Exception as ex:
                st.error(f"Error generando PDF: {ex}")

    if st.button("Generar PPT"):
        with st.spinner("Generando presentación..."):
            try:
                dr, dp, dc, ds, dl, s_r, e_r = _datos_reporte()
                ppt_bytes = generar_ppt(dr, dp, dc, s_r, e_r, df_sscc=ds, df_lim=dl)
                st.download_button(
                    "Descargar PPT",
                    data=ppt_bytes,
                    file_name=f"CTM-Reporte_{s_r}_{e_r}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as ex:
                st.error(f"Error generando PPT: {ex}")

    st.markdown(f'<p style="font-size:0.65rem;color:#94A3B8">{datetime.now().strftime("%d/%m/%Y %H:%M")}</p>', unsafe_allow_html=True)

s = fi.strftime("%Y-%m-%d")
e = ff.strftime("%Y-%m-%d")

with st.spinner("Cargando datos..."):
    df_r = load_real(s,e)
    df_p = load_prog(s,e)
    df_c = load_cmg(s, e, nodo_cmg)

if df_r.empty: st.warning("Sin datos para el período seleccionado."); st.stop()


# ── Header ────────────────────────────────────────────────────
st.markdown("# Dashboard Operacional — Complejo Térmico Mejillones")
st.markdown(f'<p style="color:#64748B;font-size:0.85rem;margin-top:-0.5rem">Período {s} → {e} · Generación real + Programada PCP + CMG {NOMBRES_NODO.get(nodo_cmg, "Crucero 220kV")}</p>', unsafe_allow_html=True)


# ── KPI cards ─────────────────────────────────────────────────
st.markdown('<div class="sec">GENERACIÓN REAL · PROMEDIO DEL PERÍODO</div>', unsafe_allow_html=True)
cols = st.columns(4)
for i,u in enumerate(["ANG1","ANG2","CCR1","CCR2"]):
    with cols[i]:
        df_u = df_r[df_r["unidad"]==u]
        if df_u.empty:
            st.markdown(f'<div class="kpi"><div class="kpi-badge" style="background:{COLORES[u]["badge"]};color:{COLORES[u]["text"]}">{LABELS[u]}</div><div class="kpi-val">—</div></div>', unsafe_allow_html=True)
            continue
        prom  = df_u["gen_real_mw"].mean()
        pmax  = PMAX.get(u, 0)
        fp    = prom/pmax*100 if pmax else 0
        df_u_sorted = df_u.sort_values("fecha_hora")
        ult_row = df_u_sorted.iloc[-1]
        ult_mw  = ult_row["gen_real_mw"]
        ult_fh  = str(ult_row["fecha_hora"])[:16]
        delta = ult_mw-prom
        sym   = "▲" if delta>=0 else "▼"
        col_d = "#10B981" if delta>=0 else "#EF4444"
        st.markdown(f"""<div class="kpi">
            <div class="kpi-badge" style="background:{COLORES[u]['badge']};color:{COLORES[u]['text']}">{LABELS[u]}</div>
            <div class="kpi-val">{prom:.1f}<span class="kpi-mw"> MW</span></div>
            <div class="kpi-sub">Factor de planta {fp:.0f}% <span style="color:#94A3B8;font-size:0.68rem">(promedio período)</span></div>
            <div class="kpi-delta" style="color:{col_d}">{sym} {abs(delta):.1f} MW vs última hora</div>
            <div style="font-size:0.68rem;color:#94A3B8;margin-top:3px">Último dato: {ult_fh}</div>
        </div>""", unsafe_allow_html=True)


# ── Función gráfico por unidad ────────────────────────────────
BG = "#FFFFFF"; GR = "#F1F5F9"

def chart_unidad(unidad: str, mostrar_desviacion: bool = False, nodo_label: str = "Crucero 220kV"):
    df_u  = df_r[df_r["unidad"]==unidad].sort_values("fecha_hora")
    df_up = df_p[df_p["unidad"]==unidad].sort_values("fecha_hora") if not df_p.empty else pd.DataFrame()
    c     = COLORES[unidad]

    tiene_cmg  = mostrar_cmg and not df_c.empty
    tiene_prog = mostrar_prog and not df_up.empty

    n_rows  = 2 if tiene_cmg else 1
    heights = [0.62, 0.38] if n_rows==2 else [1.0]

    fig = make_subplots(
        rows=n_rows, cols=1,
        shared_xaxes=True,
        row_heights=heights,
        vertical_spacing=0.12,
    )

    if df_u.empty:
        st.info(f"Sin datos para {LABELS[unidad]} en el período.")
        return

    # ── Real ──
    fig.add_trace(go.Scatter(
        x=df_u["fecha_hora"], y=df_u["gen_real_mw"],
        name="Real", mode="lines",
        line=dict(color=c["line"], width=3),
        hovertemplate="<b>Real</b> %{x|%d/%m %H:%M}<br>%{y:.1f} MW<extra></extra>",
    ), row=1, col=1)

    # ── Programada ──
    if tiene_prog:
        # Color oscurecido de la unidad para que la línea punteada sea legible
        r_d = min(255, int(c["line"][1:3],16) + 40)
        g_d = min(255, int(c["line"][3:5],16) + 30)
        b_d = min(255, int(c["line"][5:7],16) + 50)
        color_prog_dark = f"rgba({r_d},{g_d},{b_d},0.85)"
        fig.add_trace(go.Scatter(
            x=df_up["fecha_hora"], y=df_up["gen_programada_mw"],
            name="Programada", mode="lines",
            line=dict(color=color_prog_dark, width=2.2, dash="dash"),
            hovertemplate="<b>Programada</b> %{x|%d/%m %H:%M}<br>%{y:.1f} MW<extra></extra>",
        ), row=1, col=1)

        # ── Área de desviación ──
        if mostrar_desviacion:
            # Reindexar ambas series al mismo índice horario para evitar huecos
            idx_real = df_u.set_index("fecha_hora")["gen_real_mw"]
            idx_prog = df_up.set_index("fecha_hora")["gen_programada_mw"]
            idx_com  = idx_real.index.union(idx_prog.index)
            real_ri  = idx_real.reindex(idx_com).interpolate("time")
            prog_ri  = idx_prog.reindex(idx_com).interpolate("time")
            df_merge = pd.DataFrame({"real": real_ri, "prog": prog_ri}).dropna()
            if not df_merge.empty:
                r_int = int(c["line"][1:3],16)
                g_int = int(c["line"][3:5],16)
                b_int = int(c["line"][5:7],16)
                ts = df_merge.index.tolist()
                fig.add_trace(go.Scatter(
                    x=ts + ts[::-1],
                    y=df_merge["real"].tolist() + df_merge["prog"].tolist()[::-1],
                    fill="toself",
                    fillcolor=f"rgba({r_int},{g_int},{b_int},0.15)",
                    line=dict(color="rgba(0,0,0,0)"),
                    hoverinfo="skip", showlegend=True, name="Desviación",
                ), row=1, col=1)

    # ── CMG ──
    if tiene_cmg:
        r_cmg = int(c["line"][1:3],16); g_cmg = int(c["line"][3:5],16); b_cmg = int(c["line"][5:7],16)
        fig.add_trace(go.Scatter(
            x=df_c["fecha_hora"], y=df_c["cmg_usd_mwh"],
            name=f"CMG {nodo_label}", mode="lines",
            line=dict(color=c["line"], width=3),
            fill="tozeroy", fillcolor=f"rgba({r_cmg},{g_cmg},{b_cmg},0.10)",
            hovertemplate="<b>CMG</b> %{x|%d/%m %H:%M}<br>%{y:.1f} USD/MWh<extra></extra>",
        ), row=2, col=1)
        prom_cmg = df_c["cmg_usd_mwh"].mean()
        fig.add_hline(
            y=prom_cmg, line_color="#94A3B8", line_width=1, line_dash="dot",
            annotation_text=f"Prom: {prom_cmg:.1f}", annotation_position="right",
            annotation_font_color="#64748B", annotation_font_size=10,
            row=2, col=1
        )

    # ── Layout ──
    fig.update_layout(
        height=520,
        margin=dict(l=10, r=70, t=20, b=10),
        plot_bgcolor=BG, paper_bgcolor="rgba(0,0,0,0)",
        legend=dict(
            orientation="h", yanchor="bottom", y=1.02,
            xanchor="left", x=0,
            font=dict(color="#475569", size=11),
            bgcolor="rgba(0,0,0,0)",
        ),
        hovermode="x unified",
        hoverlabel=dict(bgcolor="#1E293B", font_color="#F8FAFC", bordercolor="#334155"),
    )

    # ── Eje Y ──
    fig.update_yaxes(
        title_text="MW", gridcolor=GR, zeroline=False,
        tickfont=dict(color="#94A3B8", size=10),
        title_font=dict(color="#94A3B8", size=11),
        row=1, col=1
    )
    if tiene_cmg:
        fig.update_yaxes(
            title_text="USD/MWh", gridcolor=GR, zeroline=False,
            tickfont=dict(color="#94A3B8", size=10),
            title_font=dict(color="#94A3B8", size=11),
            row=2, col=1
        )

    # ── Eje X visible en TODAS las filas ──
    for r in range(1, n_rows+1):
        fig.update_xaxes(
            showticklabels=True,
            tickformat="%d/%m\n%H:%M",
            tickfont=dict(color="#64748B", size=10),
            showgrid=False,
            row=r, col=1
        )

    st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    if mostrar_prog and df_up.empty:
        st.caption("Sin datos de programada — se importan automáticamente desde CEN PCP cada hora. El ingreso manual está disponible más abajo.")


# ── Tabs por unidad ───────────────────────────────────────────
st.markdown(f'<div class="sec">POTENCIA REAL vs PROGRAMADA + CMG {NOMBRES_NODO.get(nodo_cmg,"Crucero 220kV").upper()} · POR UNIDAD</div>', unsafe_allow_html=True)

# Checkbox desviación — solo si hay programada
hay_prog_general = not df_p.empty
mostrar_desv = False
if hay_prog_general and mostrar_prog:
    mostrar_desv = st.checkbox("Mostrar área de desviación (Real vs Programada)", value=True)

tabs = st.tabs([LABELS[u] for u in ["ANG1","ANG2","CCR1","CCR2"]])
for tab, u in zip(tabs, ["ANG1","ANG2","CCR1","CCR2"]):
    with tab:
        c = COLORES[u]
        nl = NOMBRES_NODO.get(nodo_cmg, "Crucero 220kV")
        st.markdown(f'<p style="color:{c["text"]};font-weight:600;font-size:0.9rem;margin-bottom:0.5rem">{LABELS[u]} · Real vs Programada (MW) + CMG {nl} (USD/MWh)</p>', unsafe_allow_html=True)
        chart_unidad(u, mostrar_desviacion=mostrar_desv, nodo_label=nl)


# ── Análisis de costo ─────────────────────────────────────────
st.markdown('<div class="sec">ANÁLISIS DE COSTO · CMG × GENERACIÓN</div>', unsafe_allow_html=True)

if not df_c.empty:
    df_merge = pd.merge_asof(
        df_r[["unidad","fecha_hora","gen_real_mw"]].sort_values("fecha_hora"),
        df_c[["fecha_hora","cmg_usd_mwh"]].sort_values("fecha_hora"),
        on="fecha_hora", direction="nearest", tolerance=pd.Timedelta("1h")
    )
    df_merge["ingreso_usd"] = df_merge["gen_real_mw"] * df_merge["cmg_usd_mwh"]
    ingreso_unit  = df_merge.groupby("unidad")["ingreso_usd"].sum()
    energia_unit  = df_r.groupby("unidad")["gen_real_mw"].sum()
    ingreso_total = ingreso_unit.sum()
    energia_total = energia_unit.sum()
    cmg_prom = df_c["cmg_usd_mwh"].mean()
    cmg_min  = df_c["cmg_usd_mwh"].min()
    cmg_max  = df_c["cmg_usd_mwh"].max()
    unidades_ord = ["ANG1","ANG2","CCR1","CCR2"]

    # KPIs resumen siempre visibles
    rc1,rc2,rc3,rc4 = st.columns(4)
    for col,(lbl,val,sub) in zip([rc1,rc2,rc3,rc4],[
        ("Ingreso Total Est.", f"${ingreso_total:,.0f}", "USD en el período"),
        ("Energía Total",      f"{energia_total:,.0f}",  "MWh generados"),
        ("CMG Promedio",       f"{cmg_prom:.1f}",        "USD/MWh"),
        ("Rango CMG",          f"{cmg_min:.1f} – {cmg_max:.1f}", "USD/MWh mín/máx"),
    ]):
        col.metric(lbl, val, sub)

    tab_costo_vis, tab_costo_stat = st.tabs(["Gráficos", "Estadísticas"])

    with tab_costo_vis:
        gc1, gc2 = st.columns(2)
        with gc1:
            fig2 = make_subplots(specs=[[{"secondary_y": True}]])
            fig2.add_trace(go.Bar(
                x=[LABELS[u] for u in unidades_ord],
                y=[ingreso_unit.get(u,0) for u in unidades_ord],
                name="Ingreso est. (USD)",
                marker_color=[COLORES[u]["line"] for u in unidades_ord],
                marker_opacity=0.85,
                text=[f"${ingreso_unit.get(u,0):,.0f}" for u in unidades_ord],
                textposition="outside",
                textfont=dict(size=11, color="#475569"),
            ), secondary_y=False)
            fig2.add_trace(go.Scatter(
                x=[LABELS[u] for u in unidades_ord],
                y=[energia_unit.get(u,0) for u in unidades_ord],
                name="Energía (MWh)", mode="markers+lines",
                marker=dict(size=10, color="#0F172A", symbol="diamond"),
                line=dict(color="#0F172A", width=1.5, dash="dot"),
            ), secondary_y=True)
            fig2.update_layout(
                title=dict(text="Ingreso Estimado + Energía por Unidad",
                          font=dict(size=13,color="#0F172A"), x=0),
                height=340, margin=dict(l=10, r=10, t=60, b=10),
                plot_bgcolor=BG, paper_bgcolor="rgba(0,0,0,0)",
                legend=dict(orientation="h", y=-0.15, x=0, font=dict(size=10,color="#475569")),
            )
            fig2.update_yaxes(title_text="USD", secondary_y=False,
                              gridcolor=GR, tickfont=dict(color="#94A3B8",size=10),
                              title_font=dict(color="#94A3B8",size=10))
            fig2.update_yaxes(title_text="MWh", secondary_y=True,
                              tickfont=dict(color="#94A3B8",size=10),
                              title_font=dict(color="#94A3B8",size=10), showgrid=False)
            fig2.update_xaxes(tickfont=dict(color="#475569",size=11))
            st.plotly_chart(fig2, use_container_width=True, config={"displayModeBar":False})

        with gc2:
            idx_max = df_c["cmg_usd_mwh"].idxmax()
            idx_min = df_c["cmg_usd_mwh"].idxmin()
            fig3 = go.Figure()
            fig3.add_trace(go.Scatter(
                x=df_c["fecha_hora"], y=df_c["cmg_usd_mwh"],
                mode="lines", line=dict(color=COLORES["CMG"]["line"],width=2),
                fill="tozeroy", fillcolor="rgba(109,40,217,0.12)",
                showlegend=False,
                hovertemplate="%{x|%d/%m %H:%M}<br><b>%{y:.1f} USD/MWh</b><extra></extra>",
            ))
            fig3.add_hline(y=cmg_prom, line_color="#94A3B8", line_width=1.2, line_dash="dot",
                          annotation_text=f"Prom: {cmg_prom:.1f}",
                          annotation_position="right",
                          annotation_font_color="#64748B", annotation_font_size=10)
            # halos exteriores (efecto palpitación visual)
            fig3.add_trace(go.Scatter(
                x=[df_c.loc[idx_max,"fecha_hora"]], y=[cmg_max],
                mode="markers", showlegend=False,
                marker=dict(size=22, color="rgba(239,68,68,0.18)", symbol="circle",
                            line=dict(color="rgba(239,68,68,0.45)", width=1.5)),
            ))
            fig3.add_trace(go.Scatter(
                x=[df_c.loc[idx_min,"fecha_hora"]], y=[cmg_min],
                mode="markers", showlegend=False,
                marker=dict(size=22, color="rgba(16,185,129,0.18)", symbol="circle",
                            line=dict(color="rgba(16,185,129,0.45)", width=1.5)),
            ))
            # marcadores centrales con etiqueta
            fig3.add_trace(go.Scatter(
                x=[df_c.loc[idx_max,"fecha_hora"]], y=[cmg_max],
                mode="markers+text",
                marker=dict(size=11, color="#EF4444", symbol="triangle-up",
                            line=dict(color="#fff", width=1.5)),
                text=[f" Máx: {cmg_max:.1f}"], textposition="top right",
                textfont=dict(size=10,color="#EF4444"), showlegend=False,
            ))
            fig3.add_trace(go.Scatter(
                x=[df_c.loc[idx_min,"fecha_hora"]], y=[cmg_min],
                mode="markers+text",
                marker=dict(size=11, color="#10B981", symbol="triangle-down",
                            line=dict(color="#fff", width=1.5)),
                text=[f" Mín: {cmg_min:.1f}"], textposition="bottom right",
                textfont=dict(size=10,color="#10B981"), showlegend=False,
            ))
            fig3.update_layout(
                title=dict(text="CMG en el Tiempo",
                          font=dict(size=13,color="#0F172A"), x=0),
                height=340, margin=dict(l=10, r=70, t=60, b=10),
                plot_bgcolor=BG, paper_bgcolor="rgba(0,0,0,0)",
                yaxis=dict(title="USD/MWh", gridcolor=GR,
                          tickfont=dict(color="#94A3B8",size=10),
                          title_font=dict(color="#94A3B8",size=10)),
                xaxis=dict(tickfont=dict(color="#94A3B8",size=10),
                          tickformat="%d/%m\n%H:%M", showgrid=False),
                hovermode="x unified",
                hoverlabel=dict(bgcolor="#1E293B",font_color="#F8FAFC"),
            )
            st.plotly_chart(fig3, use_container_width=True, config={"displayModeBar":False})

    with tab_costo_stat:
        BG2 = "rgba(0,0,0,0)"; GR2 = "#E2E8F0"

        # ── Fila 1: ingreso horario acumulado (área) + costo unitario por unidad ──
        df_merge_t = df_merge.copy()
        df_merge_t["fecha_hora"] = pd.to_datetime(df_merge_t["fecha_hora"])
        df_merge_t = df_merge_t.sort_values("fecha_hora")

        fig_ingh = go.Figure()
        for u in unidades_ord:
            df_u2 = df_merge_t[df_merge_t["unidad"] == u]
            fig_ingh.add_trace(go.Scatter(
                x=df_u2["fecha_hora"], y=df_u2["ingreso_usd"],
                name=LABELS[u], mode="lines",
                line=dict(color=COLORES[u]["line"], width=1.8),
                fill="tozeroy",
                fillcolor=COLORES[u]["line"].replace("#","rgba(").replace(")",",0.08)") if False else f"rgba({int(COLORES[u]['line'][1:3],16)},{int(COLORES[u]['line'][3:5],16)},{int(COLORES[u]['line'][5:7],16)},0.07)",
                hovertemplate="%{x|%d/%m %H:%M}<br><b>%{y:,.0f} USD</b><extra>" + LABELS[u] + "</extra>",
            ))
        fig_ingh.update_layout(
            title=dict(text="Ingreso estimado por hora (USD)", font=dict(size=13,color="#0F172A"), x=0),
            height=300, margin=dict(l=10,r=10,t=50,b=10),
            plot_bgcolor=BG2, paper_bgcolor=BG2,
            xaxis=dict(tickformat="%d/%m\n%H:%M", showgrid=False, tickfont=dict(color="#94A3B8",size=10)),
            yaxis=dict(gridcolor=GR2, tickfont=dict(color="#94A3B8",size=10), title="USD/h"),
            legend=dict(orientation="h", y=-0.22, font=dict(size=10)),
            hovermode="x unified",
        )

        # Costo unitario de generación (USD/MWh generado = ingreso / energía)
        costo_unit = {u: (ingreso_unit.get(u,0) / energia_unit.get(u,1)) for u in unidades_ord if energia_unit.get(u,0)>0}
        fig_cu = go.Figure(go.Bar(
            x=[LABELS[u] for u in unidades_ord if u in costo_unit],
            y=[costo_unit.get(u,0) for u in unidades_ord if u in costo_unit],
            marker_color=[COLORES[u]["line"] for u in unidades_ord if u in costo_unit],
            text=[f"{costo_unit.get(u,0):.1f}" for u in unidades_ord if u in costo_unit],
            textposition="outside",
            textfont=dict(size=11,color="#475569"),
        ))
        fig_cu.add_hline(y=cmg_prom, line_color="#94A3B8", line_dash="dot",
                         annotation_text=f"CMG prom {cmg_prom:.1f}", annotation_position="right",
                         annotation_font_color="#64748B", annotation_font_size=10)
        fig_cu.update_layout(
            title=dict(text="Ingreso medio por MWh generado (USD/MWh)", font=dict(size=13,color="#0F172A"), x=0),
            height=300, margin=dict(l=10,r=80,t=50,b=10),
            plot_bgcolor=BG2, paper_bgcolor=BG2,
            xaxis=dict(tickfont=dict(color="#475569",size=11), showgrid=False),
            yaxis=dict(gridcolor=GR2, tickfont=dict(color="#94A3B8",size=10), title="USD/MWh"),
        )

        col_s1, col_s2 = st.columns(2)
        with col_s1:
            st.plotly_chart(fig_ingh, use_container_width=True, config={"displayModeBar":False})
        with col_s2:
            st.plotly_chart(fig_cu, use_container_width=True, config={"displayModeBar":False})

        # ── Fila 2: distribución de CMG (histograma) + correlación gen vs CMG ──
        fig_hist = go.Figure(go.Histogram(
            x=df_c["cmg_usd_mwh"], nbinsx=20,
            marker_color=COLORES["CMG"]["line"], opacity=0.75,
            hovertemplate="CMG: %{x:.1f} USD/MWh<br>Frecuencia: %{y}<extra></extra>",
        ))
        fig_hist.add_vline(x=cmg_prom, line_color="#94A3B8", line_dash="dot",
                           annotation_text=f"Prom {cmg_prom:.1f}", annotation_position="top right",
                           annotation_font_color="#64748B", annotation_font_size=10)
        fig_hist.update_layout(
            title=dict(text="Distribución de precios CMG (USD/MWh)", font=dict(size=13,color="#0F172A"), x=0),
            height=300, margin=dict(l=10,r=10,t=50,b=10),
            plot_bgcolor=BG2, paper_bgcolor=BG2,
            xaxis=dict(title="USD/MWh", showgrid=False, tickfont=dict(color="#94A3B8",size=10)),
            yaxis=dict(gridcolor=GR2, tickfont=dict(color="#94A3B8",size=10), title="Horas"),
            bargap=0.08,
        )

        # Correlación generación total vs CMG
        df_corr = df_merge_t.groupby("fecha_hora").agg(
            gen_total=("gen_real_mw","sum"),
            cmg=("cmg_usd_mwh","mean"),
        ).reset_index()
        fig_corr = go.Figure(go.Scatter(
            x=df_corr["gen_total"], y=df_corr["cmg"],
            mode="markers",
            marker=dict(color=COLORES["CMG"]["line"], size=6, opacity=0.6),
            hovertemplate="Gen: %{x:.0f} MW<br>CMG: %{y:.1f} USD/MWh<extra></extra>",
        ))
        if len(df_corr) > 2:
            import numpy as np
            coef = np.polyfit(df_corr["gen_total"], df_corr["cmg"], 1)
            x_line = [df_corr["gen_total"].min(), df_corr["gen_total"].max()]
            y_line = [coef[0]*x + coef[1] for x in x_line]
            fig_corr.add_trace(go.Scatter(
                x=x_line, y=y_line, mode="lines",
                line=dict(color="#94A3B8", dash="dot", width=1.5),
                showlegend=False,
            ))
            corr_val = df_corr["gen_total"].corr(df_corr["cmg"])
            fig_corr.add_annotation(
                xref="paper", yref="paper", x=0.98, y=0.96, showarrow=False,
                text=f"r = {corr_val:.2f}",
                font=dict(size=11, color="#64748B"),
                align="right",
            )
        fig_corr.update_layout(
            title=dict(text="Generación total vs CMG (correlación horaria)", font=dict(size=13,color="#0F172A"), x=0),
            height=300, margin=dict(l=10,r=10,t=50,b=10),
            plot_bgcolor=BG2, paper_bgcolor=BG2,
            xaxis=dict(title="MW generados (total)", showgrid=False, tickfont=dict(color="#94A3B8",size=10)),
            yaxis=dict(title="CMG USD/MWh", gridcolor=GR2, tickfont=dict(color="#94A3B8",size=10)),
        )

        col_s3, col_s4 = st.columns(2)
        with col_s3:
            st.plotly_chart(fig_hist, use_container_width=True, config={"displayModeBar":False})
        with col_s4:
            st.plotly_chart(fig_corr, use_container_width=True, config={"displayModeBar":False})

        # ── Fila 3: participación de ingresos (donut) + eficiencia económica ──
        labels_pie = [LABELS[u] for u in unidades_ord if ingreso_unit.get(u,0) > 0]
        values_pie = [ingreso_unit.get(u,0) for u in unidades_ord if ingreso_unit.get(u,0) > 0]
        colors_pie = [COLORES[u]["line"] for u in unidades_ord if ingreso_unit.get(u,0) > 0]
        fig_pie = go.Figure(go.Pie(
            labels=labels_pie, values=values_pie,
            hole=0.52,
            marker=dict(colors=colors_pie),
            textfont=dict(size=11),
            hovertemplate="%{label}<br><b>$%{value:,.0f} USD</b><br>%{percent}<extra></extra>",
        ))
        fig_pie.update_layout(
            title=dict(text="Participación en ingresos estimados", font=dict(size=13,color="#0F172A"), x=0),
            height=300, margin=dict(l=10,r=10,t=50,b=10),
            paper_bgcolor=BG2,
            legend=dict(orientation="h", y=-0.1, font=dict(size=10)),
            annotations=[dict(text=f"${ingreso_total:,.0f}<br><span style='font-size:10px'>USD total</span>",
                              x=0.5, y=0.5, font_size=13, showarrow=False)],
        )

        # Eficiencia económica: ingreso por MW de capacidad instalada
        eff = {u: ingreso_unit.get(u,0) / PMAX.get(u,1) for u in unidades_ord}
        fig_eff = go.Figure(go.Bar(
            x=[LABELS[u] for u in unidades_ord],
            y=[eff[u] for u in unidades_ord],
            marker_color=[COLORES[u]["line"] for u in unidades_ord],
            text=[f"${eff[u]:,.0f}" for u in unidades_ord],
            textposition="outside",
            textfont=dict(size=11,color="#475569"),
        ))
        fig_eff.update_layout(
            title=dict(text="Ingreso estimado por MW instalado (USD/MW)", font=dict(size=13,color="#0F172A"), x=0),
            height=300, margin=dict(l=10,r=10,t=50,b=10),
            plot_bgcolor=BG2, paper_bgcolor=BG2,
            xaxis=dict(tickfont=dict(color="#475569",size=11), showgrid=False),
            yaxis=dict(gridcolor=GR2, tickfont=dict(color="#94A3B8",size=10), title="USD/MW"),
        )

        col_s5, col_s6 = st.columns(2)
        with col_s5:
            st.plotly_chart(fig_pie, use_container_width=True, config={"displayModeBar":False})
        with col_s6:
            st.plotly_chart(fig_eff, use_container_width=True, config={"displayModeBar":False})

else:
    st.info("Sin datos de CMG para calcular estadísticos de costo.")


# ── Limitaciones de Transmisión ───────────────────────────────
st.markdown('<div class="sec">LIMITACIONES DE TRANSMISIÓN — ANG / CCR</div>', unsafe_allow_html=True)

ID_UNIDAD_LABEL = {1965: "ANG1", 1966: "ANG2", 1967: "CCR1", 1968: "CCR2"}
ID_CENTRAL_LABEL = {377: "Angamos", 379: "Cochrane"}

STATUS_COLOR_LIM = {
    "pendiente":  ("#D97706", "#FEF3C7"),
    "finalizado": ("#16A34A", "#DCFCE7"),
    "anulado":    ("#94A3B8", "#F1F5F9"),
}

df_lim = load_limitaciones(s, e)

def _lim_card_html(row):
    """Devuelve el HTML de una card de limitación como string (sin llamar a st)."""
    st_key      = str(row.get("status", "")).lower()
    c_txt, c_bg = STATUS_COLOR_LIM.get(st_key, ("#475569", "#F8FAFC"))
    id_u        = int(float(row["id_unidad"])) if pd.notna(row.get("id_unidad")) else -1
    id_c        = int(float(row["id_central"])) if pd.notna(row.get("id_central")) else -1
    unidad_lbl  = ID_UNIDAD_LABEL.get(id_u, "")
    central_lbl = ID_CENTRAL_LABEL.get(id_c, str(row.get("instalacion_nombre", "")).split(" - ")[0])
    fecha_pert  = str(row.get("fecha_perturbacion") or "")[:16]
    ret_ef      = row.get("fecha_efectiva_retorno")
    ret_est     = row.get("fecha_retorno_estimada")
    es_efectivo = ret_ef and str(ret_ef) not in ("NaT", "None", "nan", "")
    retorno_val = ret_ef if es_efectivo else ret_est
    fecha_ret   = str(retorno_val)[:16] if (retorno_val and str(retorno_val) not in ("NaT", "None", "nan", "")) else "—"
    ret_label   = "Cierre real" if es_efectivo else "Retorno est."
    potencia    = row.get("potencia")
    pot_str     = str(int(float(potencia))) + " " + str(row.get("unidad_medida_potencia") or "MW") if pd.notna(potencia) and float(potencia) > 0 else ""
    correlativo = row.get("correlativo")
    corr_num    = str(int(float(correlativo))) if pd.notna(correlativo) else ""
    obs         = str(row.get("observacion") or "").strip()[:220]
    afecta      = bool(row.get("afecta_sscc"))

    partes_f1 = []
    _badge_cls = ' class="badge-pend"' if st_key == "pendiente" else ''
    partes_f1.append(f'<span{_badge_cls} style="background:{c_bg};color:{c_txt};padding:2px 8px;border-radius:4px;font-size:0.71rem;font-weight:700;text-transform:uppercase">{row.get("status","")}</span>')
    partes_f1.append(f'<span style="font-weight:600;font-size:0.84rem">{central_lbl}</span>')
    if unidad_lbl:
        partes_f1.append(f'<span style="background:#EDE9FE;color:#6D28D9;padding:1px 7px;border-radius:4px;font-size:0.71rem;font-weight:600">{unidad_lbl}</span>')
    if afecta:
        partes_f1.append('<span style="background:#FEF3C7;color:#D97706;padding:1px 6px;border-radius:4px;font-size:0.67rem">Afecta SSCC</span>')
    if corr_num:
        partes_f1.append(f'<span style="font-size:0.67rem;color:#94A3B8">N. {corr_num}</span>')

    f2_left  = f'<span style="font-size:0.71rem;color:#475569">Apertura: <b>{fecha_pert}</b> &rarr; {ret_label}: <b>{fecha_ret}</b></span>'
    f2_right = f'<span style="font-size:0.77rem;font-weight:700;color:#DC2626">{pot_str}</span>' if pot_str else ""
    fila2    = f'{f2_left}{("&nbsp;&nbsp;&nbsp;" + f2_right) if f2_right else ""}'
    obs_div  = f'<div style="font-size:0.71rem;color:#64748B;margin-top:3px">{obs}</div>' if obs else ""

    return (
        f'<div style="border:1px solid #E2E8F0;border-radius:8px;padding:10px 14px;margin-bottom:8px;background:#FAFAFA">'
        f'<div style="display:flex;gap:6px;flex-wrap:wrap;align-items:center;margin-bottom:3px">{"".join(partes_f1)}</div>'
        f'<div>{fila2}</div>'
        f'{obs_div}</div>'
    )

if df_lim.empty:
    st.info("Sin limitaciones registradas para el período seleccionado.")
else:
    activas   = df_lim[df_lim["status"] == "pendiente"]
    n_activas = len(activas)
    n_total   = len(df_lim)
    n_sscc    = int(df_lim["afecta_sscc"].fillna(False).sum())
    pot_max   = df_lim[df_lim["status"] == "pendiente"]["potencia"].max()
    pot_str   = f"{pot_max:.0f} MW" if pd.notna(pot_max) and pot_max > 0 else "—"

    kl1, kl2, kl3, kl4 = st.columns(4)
    kl1.metric("Activas (pendiente)",    n_activas)
    kl2.metric("Total en período",        n_total)
    kl3.metric("Afectan SSCC",           n_sscc)
    kl4.metric("Mayor limitación activa", pot_str)

    df_lim["_unidad"] = df_lim["id_unidad"].apply(
        lambda x: ID_UNIDAD_LABEL.get(int(float(x)), "") if pd.notna(x) else ""
    )
    df_lim_sorted = df_lim.sort_values("fecha_perturbacion", ascending=False)
    MAX_CARDS = 5

    tabs_lim = st.tabs(["ANG1", "ANG2", "CCR1", "CCR2", "Todas", "Estadísticas"])
    for tab, unidad in zip(tabs_lim[:4], ["ANG1", "ANG2", "CCR1", "CCR2"]):
        with tab:
            df_u = df_lim_sorted[df_lim_sorted["_unidad"] == unidad]
            if df_u.empty:
                st.info(f"Sin limitaciones para {unidad} en el período.")
            else:
                df_u_show = df_u.head(MAX_CARDS)
                cards_html = "".join(_lim_card_html(row) for _, row in df_u_show.iterrows())
                st.markdown(cards_html, unsafe_allow_html=True)
                if len(df_u) > MAX_CARDS:
                    st.caption(f"+{len(df_u) - MAX_CARDS} más en «Todas»")
    with tabs_lim[4]:
        cards_html = "".join(_lim_card_html(row) for _, row in df_lim_sorted.iterrows())
        st.markdown(cards_html, unsafe_allow_html=True)

    with tabs_lim[5]:
        df_stat = df_lim.copy()
        df_stat["fecha_perturbacion"] = pd.to_datetime(df_stat["fecha_perturbacion"])

        # ── Fila 1: barras por mes + donut por unidad ─────────────
        gc1, gc2 = st.columns([3, 2])

        with gc1:
            df_mes = df_stat.copy()
            df_mes["mes"] = df_mes["fecha_perturbacion"].dt.to_period("M").astype(str)
            pivot_mes = (
                df_mes.groupby(["mes", "status"])
                .size()
                .reset_index(name="n")
            )
            COLOR_STATUS = {"pendiente": "#D97706", "finalizado": "#16A34A", "anulado": "#94A3B8"}
            fig_mes = go.Figure()
            for st_val in ["pendiente", "finalizado", "anulado"]:
                d = pivot_mes[pivot_mes["status"] == st_val]
                if not d.empty:
                    fig_mes.add_trace(go.Bar(
                        x=d["mes"], y=d["n"],
                        name=st_val.capitalize(),
                        marker_color=COLOR_STATUS[st_val],
                    ))
            fig_mes.update_layout(
                title="Limitaciones por mes",
                barmode="stack",
                height=300,
                margin=dict(t=40, b=30, l=30, r=10),
                legend=dict(orientation="h", y=-0.25),
                plot_bgcolor="white",
                yaxis=dict(gridcolor="#F1F5F9"),
            )
            fig_mes.update_xaxes(tickangle=-30)
            st.plotly_chart(fig_mes, use_container_width=True)

        with gc2:
            COLORES_UNIDAD = {"ANG1": "#7C3AED", "ANG2": "#2563EB", "CCR1": "#CA8A04", "CCR2": "#16A34A"}
            conteo_u = df_stat["_unidad"].replace("", "Sin unidad").value_counts().reset_index()
            conteo_u.columns = ["unidad", "n"]
            colores_donut = [COLORES_UNIDAD.get(u, "#94A3B8") for u in conteo_u["unidad"]]
            fig_donut = go.Figure(go.Pie(
                labels=conteo_u["unidad"],
                values=conteo_u["n"],
                hole=0.55,
                marker_colors=colores_donut,
                textinfo="label+value",
                textfont_size=12,
            ))
            fig_donut.update_layout(
                title="Por unidad",
                height=300,
                margin=dict(t=40, b=10, l=10, r=10),
                showlegend=False,
            )
            st.plotly_chart(fig_donut, use_container_width=True)

        # ── Fila 2: histograma potencia + duración finalizadas ─────
        gc3, gc4 = st.columns(2)

        with gc3:
            df_pot = df_stat[df_stat["potencia"].notna() & (df_stat["potencia"] > 0)].copy()
            bins = [0, 50, 100, 150, 200, 300]
            labels_pot = ["1–50", "51–100", "101–150", "151–200", ">200"]
            df_pot["rango"] = pd.cut(df_pot["potencia"], bins=bins, labels=labels_pot, right=True)
            conteo_pot = df_pot["rango"].value_counts().reindex(labels_pot, fill_value=0).reset_index()
            conteo_pot.columns = ["rango", "n"]
            fig_pot = go.Figure(go.Bar(
                x=conteo_pot["rango"],
                y=conteo_pot["n"],
                marker_color="#3B82F6",
                text=conteo_pot["n"],
                textposition="outside",
            ))
            fig_pot.update_layout(
                title="Distribución por potencia limitada (MW)",
                height=300,
                margin=dict(t=40, b=30, l=30, r=10),
                plot_bgcolor="white",
                yaxis=dict(gridcolor="#F1F5F9", title="Cantidad"),
                xaxis=dict(title="Rango MW"),
            )
            st.plotly_chart(fig_pot, use_container_width=True)

        with gc4:
            df_fin = df_stat[df_stat["status"] == "finalizado"].copy()
            df_fin["retorno"] = pd.to_datetime(df_fin["fecha_efectiva_retorno"].fillna(df_fin["fecha_retorno_estimada"]))
            df_fin["duracion_dias"] = (df_fin["retorno"] - df_fin["fecha_perturbacion"]).dt.total_seconds() / 86400
            df_fin = df_fin[df_fin["duracion_dias"].notna() & (df_fin["duracion_dias"] >= 0)].sort_values("fecha_perturbacion")
            df_fin["label"] = df_fin["correlativo"].apply(lambda x: f"N.{int(float(x))}" if pd.notna(x) else "")
            if df_fin.empty:
                st.info("Sin limitaciones finalizadas en el período para calcular duración.")
            else:
                fig_dur = go.Figure(go.Bar(
                    x=df_fin["label"],
                    y=df_fin["duracion_dias"].round(1),
                    marker_color="#16A34A",
                    text=df_fin["duracion_dias"].round(1),
                    textposition="outside",
                ))
                fig_dur.update_layout(
                    title="Duración (días) — limitaciones finalizadas",
                    height=300,
                    margin=dict(t=40, b=30, l=30, r=10),
                    plot_bgcolor="white",
                    yaxis=dict(gridcolor="#F1F5F9", title="Días"),
                    xaxis=dict(title="Correlativo"),
                )
                st.plotly_chart(fig_dur, use_container_width=True)

    # Tabla completa con <details>/<summary> HTML nativo (evita bug de st.expander post-tabs)
    df_lim_sorted2 = df_lim_sorted.copy()
    df_lim_sorted2["correlativo"] = df_lim_sorted2["correlativo"].apply(
        lambda x: str(int(float(x))) if pd.notna(x) else ""
    )
    cols_tabla = ["correlativo", "status", "instalacion_nombre", "fecha_perturbacion",
                  "fecha_retorno_estimada", "fecha_efectiva_retorno", "potencia", "afecta_sscc", "observacion"]
    df_t = df_lim_sorted2[[c for c in cols_tabla if c in df_lim_sorted2.columns]]
    hdrs = "".join(f'<th style="padding:6px 10px;text-align:left;font-size:0.72rem;color:#475569;border-bottom:1px solid #E2E8F0;white-space:nowrap">{c}</th>' for c in df_t.columns)
    rows_html = ""
    for _, r in df_t.iterrows():
        cells = "".join(f'<td style="padding:5px 10px;font-size:0.72rem;border-bottom:1px solid #F1F5F9;white-space:nowrap">{str(r[c]) if pd.notna(r[c]) else ""}</td>' for c in df_t.columns)
        rows_html += f"<tr>{cells}</tr>"
    tabla_html = (
        f'<details style="background:#F8FAFC;border:1px solid #E2E8F0;border-radius:8px;padding:0.6rem 1rem;margin-top:0.5rem">'
        f'<summary style="cursor:pointer;font-weight:600;color:#334155;font-size:0.88rem">Ver tabla completa de limitaciones ({len(df_t)} registros)</summary>'
        f'<div style="overflow-x:auto;margin-top:0.6rem"><table style="border-collapse:collapse;width:100%"><thead><tr>{hdrs}</tr></thead><tbody>{rows_html}</tbody></table></div>'
        f'</details>'
    )
    st.markdown(tabla_html, unsafe_allow_html=True)



# ── Servicios Complementarios (SSCC) ─────────────────────────
st.markdown('<div class="sec">SERVICIOS COMPLEMENTARIOS (SSCC)</div>', unsafe_allow_html=True)

COLORES_SSCC = {
    "CSF(+)": "#16A34A", "CSF(-)": "#0891B2",
    "CPF(+)": "#6D28D9", "CPF(-)": "#D97706",
    "CT":     "#CA8A04",
    "CTF":    "#DC2626",
}
BADGE_SSCC = {
    "CSF(+)": "#DCFCE7", "CSF(-)": "#CFFAFE",
    "CPF(+)": "#EDE9FE", "CPF(-)": "#FEF3C7",
    "CT":     "#FEF9C3",
    "CTF":    "#FEE2E2",
}
st.markdown("""
<details style="background:#F8FAFC;border:1px solid #E2E8F0;border-radius:8px;padding:0.7rem 1rem;margin-bottom:1rem;">
<summary style="cursor:pointer;font-weight:600;color:#334155;font-size:0.9rem;">Guía de instrucciones SSCC — CSF, CPF, CT, CTF</summary>
<div style="margin-top:0.8rem;font-size:0.88rem;color:#475569;line-height:1.6;">
<p>Los <strong>Servicios Complementarios (SSCC)</strong> son prestaciones que el Coordinador Eléctrico Nacional instruye
a las unidades generadoras para mantener la seguridad y calidad del Sistema Eléctrico Nacional,
más allá de su generación de energía.</p>
<table style="width:100%;border-collapse:collapse;margin:0.5rem 0;">
<thead><tr style="background:#E2E8F0;">
<th style="padding:6px 12px;text-align:left;font-size:0.82rem;">Instrucción</th>
<th style="padding:6px 12px;text-align:left;font-size:0.82rem;">Significado</th>
</tr></thead>
<tbody>
<tr><td style="padding:5px 12px;border-top:1px solid #E2E8F0;"><strong>CSF(+)</strong></td><td style="padding:5px 12px;border-top:1px solid #E2E8F0;">Control Secundario de Frecuencia en <strong>subida</strong> — la unidad debe estar disponible para aumentar potencia y corregir la frecuencia del sistema</td></tr>
<tr><td style="padding:5px 12px;border-top:1px solid #E2E8F0;"><strong>CSF(−)</strong></td><td style="padding:5px 12px;border-top:1px solid #E2E8F0;">Control Secundario de Frecuencia en <strong>bajada</strong> — la unidad debe estar disponible para reducir potencia</td></tr>
<tr><td style="padding:5px 12px;border-top:1px solid #E2E8F0;"><strong>CPF(+)</strong></td><td style="padding:5px 12px;border-top:1px solid #E2E8F0;">Control Primario de Frecuencia en <strong>subida</strong> — respuesta automática e inmediata ante caída de frecuencia</td></tr>
<tr><td style="padding:5px 12px;border-top:1px solid #E2E8F0;"><strong>CPF(−)</strong></td><td style="padding:5px 12px;border-top:1px solid #E2E8F0;">Control Primario de Frecuencia en <strong>bajada</strong> — respuesta automática ante alza de frecuencia</td></tr>
<tr><td style="padding:5px 12px;border-top:1px solid #E2E8F0;"><strong>CT</strong></td><td style="padding:5px 12px;border-top:1px solid #E2E8F0;">Control de <strong>Tensión</strong> — regulación de tensión reactiva en la barra de conexión</td></tr>
<tr><td style="padding:5px 12px;border-top:1px solid #E2E8F0;"><strong>CTF</strong></td><td style="padding:5px 12px;border-top:1px solid #E2E8F0;">Control <strong>Terciario de Frecuencia</strong> — reserva de potencia de respuesta más lenta (minutos) que activa el Coordinador para restablecer la frecuencia nominal tras un evento, liberando la reserva secundaria (CSF) para nuevas contingencias</td></tr>
</tbody>
</table>
<p>Cada instrucción indica un <strong>período de prestación</strong> (inicio → fin) durante el cual la unidad debe
mantener una reserva de potencia disponible para activación. El campo <strong>Disp. MW</strong> indica la capacidad
comprometida cuando está declarada.</p>
</div>
</details>
""", unsafe_allow_html=True)

df_sscc = load_sscc(s, e)

if df_sscc.empty:
    st.info("Sin datos SSCC para el período seleccionado. Los datos se adquieren automáticamente cada hora.")
else:
    total_instr  = len(df_sscc)
    unidades_act = df_sscc["unidad"].nunique()
    tipos_act    = df_sscc["instruccion_sscc"].nunique()
    dias_con_dat = df_sscc["fecha"].nunique()

    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Instrucciones totales", total_instr)
    k2.metric("Unidades con SSCC",     f"{unidades_act} / 4")
    k3.metric("Tipos de servicio",     tipos_act)
    k4.metric("Días con datos",        dias_con_dat)

    st.markdown("")

    tab_s1, tab_s2, tab_s3 = st.tabs(["Por unidad", "Estadísticas", "Tabla completa"])

    with tab_s1:
        cols_unit = st.columns(4)
        for col, unidad in zip(cols_unit, ["ANG1","ANG2","CCR1","CCR2"]):
            df_u = df_sscc[df_sscc["unidad"] == unidad].sort_values(
                ["fecha", "inicio_periodo"], ascending=[False, False]
            )
            total_u = len(df_u)
            df_u_show = df_u.head(5)
            with col:
                st.markdown(f"**{LABELS[unidad]}**")
                if df_u.empty:
                    st.caption("Sin instrucciones")
                else:
                    for idx_row, (_, row) in enumerate(df_u_show.iterrows()):
                        tipo  = str(row["instruccion_sscc"])
                        color = COLORES_SSCC.get(tipo, "#64748B")
                        bg    = BADGE_SSCC.get(tipo, "#F1F5F9")
                        ini   = str(row["inicio_periodo"])[:5] if row["inicio_periodo"] else "—"
                        fin   = str(row["fin_periodo"])[:5]    if row["fin_periodo"]    else "—"
                        fecha = str(row["fecha"])
                        extra_cls = ' sscc-latest' if idx_row == 0 else ''
                        st.markdown(
                            f'<div class="sscc-card{extra_cls}" style="border:1px solid {color}33;border-left:3px solid {color};'
                            f'background:{bg};border-radius:6px;padding:6px 10px;margin-bottom:6px;">'
                            f'<span style="font-weight:700;color:{color};font-size:0.8rem">{tipo}</span>'
                            f'<span style="color:#64748B;font-size:0.72rem;float:right">{fecha}</span><br>'
                            f'<span style="color:#475569;font-size:0.72rem">{ini} → {fin}</span>'
                            f'</div>',
                            unsafe_allow_html=True,
                        )
                    if total_u > 5:
                        st.caption(f"+{total_u - 5} más en «Tabla completa»")

    with tab_s2:
        BG = "rgba(0,0,0,0)"; GR = "#E2E8F0"

        # Gráfico 1: instrucciones por tipo (barras)
        conteo_tipo = df_sscc.groupby("instruccion_sscc").size().reset_index(name="count")
        conteo_tipo = conteo_tipo.sort_values("count", ascending=False)
        fig_tipos = go.Figure(go.Bar(
            x=conteo_tipo["instruccion_sscc"],
            y=conteo_tipo["count"],
            marker_color=[COLORES_SSCC.get(t, "#64748B") for t in conteo_tipo["instruccion_sscc"]],
            text=conteo_tipo["count"], textposition="outside",
        ))
        fig_tipos.update_layout(
            title=dict(text="Instrucciones por tipo de SSCC", font=dict(size=13, color="#0F172A"), x=0),
            height=280, margin=dict(l=10, r=10, t=50, b=10),
            plot_bgcolor=BG, paper_bgcolor=BG,
            xaxis=dict(tickfont=dict(color="#94A3B8", size=11), showgrid=False),
            yaxis=dict(gridcolor=GR, tickfont=dict(color="#94A3B8", size=10), title="N° instrucciones"),
        )

        # Gráfico 2: instrucciones por unidad (barras apiladas por tipo)
        pivot_unidad = df_sscc.groupby(["unidad","instruccion_sscc"]).size().reset_index(name="count")
        fig_unidad = go.Figure()
        for tipo in df_sscc["instruccion_sscc"].unique():
            df_t = pivot_unidad[pivot_unidad["instruccion_sscc"] == tipo]
            fig_unidad.add_trace(go.Bar(
                name=tipo,
                x=df_t["unidad"],
                y=df_t["count"],
                marker_color=COLORES_SSCC.get(tipo, "#64748B"),
            ))
        fig_unidad.update_layout(
            barmode="stack",
            title=dict(text="Instrucciones por unidad", font=dict(size=13, color="#0F172A"), x=0),
            height=280, margin=dict(l=10, r=10, t=50, b=10),
            plot_bgcolor=BG, paper_bgcolor=BG,
            xaxis=dict(tickfont=dict(color="#94A3B8", size=11), showgrid=False),
            yaxis=dict(gridcolor=GR, tickfont=dict(color="#94A3B8", size=10), title="N° instrucciones"),
            legend=dict(orientation="h", y=-0.2, font=dict(size=10)),
        )

        col_g1, col_g2 = st.columns(2)
        with col_g1:
            st.plotly_chart(fig_tipos, use_container_width=True, config={"displayModeBar": False})
        with col_g2:
            st.plotly_chart(fig_unidad, use_container_width=True, config={"displayModeBar": False})

        # Gráfico 3: duración promedio por tipo (horas)
        df_dur = df_sscc.copy()
        def parse_hhmm(t):
            try:
                h, m, *_ = str(t).split(":")
                return int(h) + int(m) / 60
            except:
                return None
        df_dur["h_ini"] = df_dur["inicio_periodo"].apply(parse_hhmm)
        df_dur["h_fin"] = df_dur["fin_periodo"].apply(parse_hhmm)
        df_dur["duracion_h"] = (df_dur["h_fin"] - df_dur["h_ini"]).clip(lower=0)
        dur_tipo = df_dur.groupby("instruccion_sscc")["duracion_h"].mean().reset_index()
        dur_tipo = dur_tipo.sort_values("duracion_h", ascending=False)

        fig_dur = go.Figure(go.Bar(
            x=dur_tipo["instruccion_sscc"],
            y=dur_tipo["duracion_h"].round(1),
            marker_color=[COLORES_SSCC.get(t, "#64748B") for t in dur_tipo["instruccion_sscc"]],
            text=dur_tipo["duracion_h"].round(1).astype(str) + " h",
            textposition="outside",
        ))
        fig_dur.update_layout(
            title=dict(text="Duración promedio por tipo (horas)", font=dict(size=13, color="#0F172A"), x=0),
            height=280, margin=dict(l=10, r=10, t=50, b=10),
            plot_bgcolor=BG, paper_bgcolor=BG,
            xaxis=dict(tickfont=dict(color="#94A3B8", size=11), showgrid=False),
            yaxis=dict(gridcolor=GR, tickfont=dict(color="#94A3B8", size=10), title="Horas promedio"),
        )

        # Gráfico 4: evolución diaria de instrucciones
        if dias_con_dat > 1:
            evol = df_sscc.groupby(["fecha","instruccion_sscc"]).size().reset_index(name="count")
            evol["fecha"] = pd.to_datetime(evol["fecha"])
            fig_evol = go.Figure()
            for tipo in sorted(df_sscc["instruccion_sscc"].unique()):
                df_e = evol[evol["instruccion_sscc"] == tipo]
                fig_evol.add_trace(go.Scatter(
                    x=df_e["fecha"], y=df_e["count"], name=tipo,
                    mode="lines+markers",
                    line=dict(color=COLORES_SSCC.get(tipo, "#64748B"), width=2),
                    marker=dict(size=6),
                ))
            fig_evol.update_layout(
                title=dict(text="Evolución diaria de instrucciones SSCC", font=dict(size=13, color="#0F172A"), x=0),
                height=280, margin=dict(l=10, r=10, t=50, b=10),
                plot_bgcolor=BG, paper_bgcolor=BG,
                xaxis=dict(tickfont=dict(color="#94A3B8", size=10), showgrid=False, tickformat="%d/%m"),
                yaxis=dict(gridcolor=GR, tickfont=dict(color="#94A3B8", size=10), title="N° instrucciones"),
                legend=dict(orientation="h", y=-0.2, font=dict(size=10)),
                hovermode="x unified",
            )
            col_g3, col_g4 = st.columns(2)
            with col_g3:
                st.plotly_chart(fig_dur, use_container_width=True, config={"displayModeBar": False})
            with col_g4:
                st.plotly_chart(fig_evol, use_container_width=True, config={"displayModeBar": False})
        else:
            st.plotly_chart(fig_dur, use_container_width=True, config={"displayModeBar": False})

    with tab_s3:
        df_show_sscc = df_sscc.copy()
        df_show_sscc["inicio"] = df_show_sscc["inicio_periodo"].astype(str).str[:5]
        df_show_sscc["fin"]    = df_show_sscc["fin_periodo"].astype(str).str[:5]
        df_show_sscc["motivo"] = df_show_sscc["motivo"].fillna("").str[:80]
        st.dataframe(
            df_show_sscc[["fecha","unidad","instruccion_sscc","inicio","fin","disponibilidad","motivo","estado_sabana"]].rename(columns={
                "instruccion_sscc": "Instrucción",
                "inicio":           "Inicio",
                "fin":              "Fin",
                "disponibilidad":   "Disp. MW",
                "motivo":           "Motivo",
                "estado_sabana":    "Estado",
            }),
            use_container_width=True, hide_index=True,
        )


# ── Solicitudes de Trabajo ────────────────────────────────────
st.markdown('<div class="sec">SOLICITUDES DE TRABAJO — AES ANDES / ANG / CCR</div>', unsafe_allow_html=True)

STATUS_COLOR_SOL = {
    "pendiente":          ("#D97706", "#FEF3C7"),
    "ejecucion_exitosa":  ("#16A34A", "#DCFCE7"),
    "anulado":            ("#94A3B8", "#F1F5F9"),
    "borrador":           ("#6366F1", "#EEF2FF"),
}
TIPO_LABEL = {"desconexion": "Desconexión", "intervencion": "Intervención"}
TYPE_LABEL = {"central_generadora": "Central", "subestacion": "Subestación", "linea": "Línea"}

df_sol = load_solicitudes(s, e)

if df_sol.empty:
    st.info("Sin solicitudes de trabajo registradas para el período seleccionado.")
else:
    # KPIs
    n_total    = len(df_sol)
    n_pend     = (df_sol["status"] == "pendiente").sum()
    n_ejec     = (df_sol["status"] == "ejecucion_exitosa").sum()
    n_desc     = (df_sol["tipo_solicitud"] == "desconexion").sum()

    kc1, kc2, kc3, kc4 = st.columns(4)
    kc1.metric("Total solicitudes", n_total)
    kc2.metric("Pendientes", int(n_pend))
    kc3.metric("Ejecutadas", int(n_ejec))
    kc4.metric("Desconexiones", int(n_desc))

    # Tabs
    tab_s_todas, tab_s_pend, tab_s_tabla = st.tabs(["Todas", "Pendientes", "Tabla completa"])

    def _sol_cards(df_view):
        if df_view.empty:
            st.info("Sin registros."); return
        for _, row in df_view.head(5).iterrows():
            st_key      = str(row.get("status", "")).lower()
            c_txt, c_bg = STATUS_COLOR_SOL.get(st_key, ("#475569", "#F8FAFC"))
            tipo_lbl    = TIPO_LABEL.get(str(row.get("tipo_solicitud","")), str(row.get("tipo_solicitud","")))
            type_lbl    = TYPE_LABEL.get(str(row.get("type","")), str(row.get("type","")))
            inst        = str(row.get("instalacion_nombre") or "—")
            f_ini       = str(row.get("fecha_inicio") or "—")[:16]
            f_fin       = str(row.get("fecha_fin")    or "—")[:16]
            corr        = int(row["correlativo"]) if pd.notna(row.get("correlativo")) else "—"
            riesgo      = str(row.get("descripcion_nivel_riesgo") or "—")
            st.markdown(f"""
            <div style="background:{c_bg};border-left:4px solid {c_txt};border-radius:6px;
                        padding:0.6rem 0.9rem;margin-bottom:0.5rem;font-size:0.82rem">
                <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:0.3rem">
                    <span style="font-weight:700;color:#1E293B">N° {corr} — {inst}</span>
                    <span style="background:{c_txt};color:#fff;border-radius:4px;
                                 padding:2px 8px;font-size:0.72rem;font-weight:600">
                        {st_key.replace("_"," ").upper()}
                    </span>
                </div>
                <div style="color:#475569;line-height:1.7">
                    <b>Tipo:</b> {tipo_lbl} · {type_lbl} &nbsp;|&nbsp;
                    <b>Inicio:</b> {f_ini} &nbsp;|&nbsp; <b>Fin:</b> {f_fin}<br>
                    <b>Riesgo:</b> {riesgo[:120]}{'…' if len(riesgo)>120 else ''}
                </div>
            </div>
            """, unsafe_allow_html=True)
        if len(df_view) > 5:
            st.caption(f"+{len(df_view)-5} más en «Tabla completa»")

    with tab_s_todas:
        _sol_cards(df_sol)

    with tab_s_pend:
        df_pend = df_sol[df_sol["status"] == "pendiente"]
        _sol_cards(df_pend)

    with tab_s_tabla:
        cols_show = ["correlativo", "empresa_nombre", "instalacion_nombre", "status",
                     "tipo_solicitud", "type", "fecha_inicio", "fecha_fin"]
        df_t = df_sol[cols_show].copy()
        df_t.columns = ["Correlativo", "Empresa", "Instalación", "Status",
                        "Tipo", "Elemento", "Inicio", "Fin"]
        hdrs = "".join(f'<th style="padding:5px 10px;background:#F1F5F9;font-size:0.72rem;text-align:left">{c}</th>' for c in df_t.columns)
        rows_html = ""
        for _, r in df_t.iterrows():
            cells = "".join(f'<td style="padding:5px 10px;font-size:0.72rem;border-bottom:1px solid #F1F5F9;white-space:nowrap">{str(r[c]) if pd.notna(r[c]) else ""}</td>' for c in df_t.columns)
            rows_html += f"<tr>{cells}</tr>"
        tabla_html = (
            f'<div style="overflow-x:auto;margin-top:0.5rem">'
            f'<table style="border-collapse:collapse;width:100%">'
            f'<thead><tr>{hdrs}</tr></thead><tbody>{rows_html}</tbody></table></div>'
        )
        st.markdown(tabla_html, unsafe_allow_html=True)


# ── Potencia programada ───────────────────────────────────────
st.markdown('<div class="sec">POTENCIA PROGRAMADA · CEN PCP + INGRESO MANUAL</div>', unsafe_allow_html=True)
st.info(
    "Los datos de programación se importan automáticamente desde la API CEN (PCP) cada hora. "
    "El ingreso manual permite agregar o corregir valores de respaldo cuando no hay datos PCP disponibles.",
)

tab_p1, tab_p2 = st.tabs(["Por hora", "24 horas de una vez"])

with tab_p1:
    u_prog = st.radio("Unidad", ["ANG1","ANG2","CCR1","CCR2"], key="up", horizontal=True)
    pc2,pc3,pc4 = st.columns(3)
    with pc2: f_prog  = st.date_input("Fecha",value=hoy,max_value=hoy,key="fp")
    with pc3: h_prog  = st.number_input("Hora (1-24)",1,24,datetime.now().hour+1,key="hp")
    with pc4: mw_prog = st.number_input("MW programados",0.0,400.0,step=0.5,key="mwp")
    if st.button("Guardar hora",key="btn_h"):
        fh = f"{f_prog} {int(h_prog)-1:02d}:00:00"
        ok = exe("""INSERT INTO generacion_programada (unidad,gen_programada_mw,fecha_hora,hora,fuente)
                    VALUES (%s,%s,%s,%s,'MANUAL')
                    ON CONFLICT (unidad,fecha_hora,fuente) DO UPDATE SET gen_programada_mw=EXCLUDED.gen_programada_mw""",
                 (u_prog,mw_prog,fh,int(h_prog)))
        if ok: st.success(f"Guardado: {LABELS[u_prog]} H{h_prog} → {mw_prog} MW"); st.cache_data.clear(); st.rerun()

with tab_p2:
    st.caption("Selecciona unidad y fecha, luego pega los 24 valores MW separados por salto de línea (hora 1 → hora 24).")
    u_masa = st.radio("Unidad", ["ANG1","ANG2","CCR1","CCR2"], key="um", horizontal=True)
    mc1,mc2 = st.columns([1,2])
    with mc1:
        f_masa = st.date_input("Fecha",value=hoy,max_value=hoy,key="fm")
        st.caption("Ejemplo formato:\n280.5\n275.3\n271.0\n...")
    with mc2:
        mw_masa = st.text_area("24 valores MW (uno por línea):", height=220, key="mwm",
                               placeholder="280.5\n275.3\n271.0\n268.5\n...")
    if st.button("Guardar las 24 horas",key="btn_m"):
        try:
            valores = [float(v.strip().replace(",",".")) for v in mw_masa.strip().split("\n") if v.strip()]
            if len(valores) != 24:
                st.error(f"Se esperan 24 valores, se ingresaron {len(valores)}.")
            else:
                params_list = [(u_masa, mw, f"{f_masa} {h-1:02d}:00:00", h, "MANUAL") for h,mw in enumerate(valores,1)]
                ok = exe_many("""INSERT INTO generacion_programada (unidad,gen_programada_mw,fecha_hora,hora,fuente)
                                 VALUES (%s,%s,%s,%s,%s)
                                 ON CONFLICT (unidad,fecha_hora,fuente) DO UPDATE SET gen_programada_mw=EXCLUDED.gen_programada_mw""",
                              params_list)
                if ok: st.success(f"24 horas guardadas: {LABELS[u_masa]} · {f_masa}"); st.cache_data.clear(); st.rerun()
        except ValueError:
            st.error("Formato inválido. Solo números, uno por línea.")

df_pv2 = load_prog(s,e)
if not df_pv2.empty:
    show_prog_table = st.checkbox("Ver tabla de datos programados", value=False, key="show_prog_tbl")
    if show_prog_table:
        df_show = df_pv2.copy()
        df_show["fecha_hora"] = pd.to_datetime(df_show["fecha_hora"]).dt.strftime("%Y-%m-%d %H:%M")
        cols_show = ["unidad","fecha_hora","hora","gen_programada_mw","fuente"] if "fuente" in df_show.columns else ["unidad","fecha_hora","hora","gen_programada_mw"]
        st.dataframe(df_show[cols_show].rename(
            columns={"gen_programada_mw":"MW Programado","fuente":"Fuente"}), use_container_width=True, hide_index=True)

    st.markdown("---")
    st.markdown("**Modificar o eliminar registro programada:**")
    # Necesitamos id — hacer query con id
    df_prog_crud = qry("SELECT id,unidad,gen_programada_mw,fecha_hora,hora FROM generacion_programada WHERE fecha_hora::date BETWEEN %s AND %s AND fuente='MANUAL' ORDER BY unidad,fecha_hora",(s,e))
    if not df_prog_crud.empty:
        df_prog_crud["fecha_hora"] = pd.to_datetime(df_prog_crud["fecha_hora"])
        opc_p = df_prog_crud.apply(lambda r: f"[{r['unidad']}] {r['fecha_hora'].strftime('%d/%m %H:%M')} — {r['gen_programada_mw']:.1f} MW", axis=1).tolist()
        idx_p = st.selectbox("Seleccionar registro programada", range(len(opc_p)),
                             format_func=lambda i: opc_p[i], key="sel_prog")
        reg_p = df_prog_crud.iloc[idx_p]
        col_ep, col_dp = st.columns([2,1])
        with col_ep:
            new_mw_p = st.number_input("Nuevo valor MW:", value=float(reg_p["gen_programada_mw"]),
                                       min_value=0.0, max_value=400.0, step=0.5, key="edit_mw_p")
            if st.button("Actualizar MW", key="upd_prog"):
                ok = exe("UPDATE generacion_programada SET gen_programada_mw=%s WHERE id=%s",
                         (new_mw_p, int(reg_p["id"])))
                if ok: st.success("Actualizado."); st.cache_data.clear(); st.rerun()
        with col_dp:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Eliminar", key="del_prog", type="primary"):
                ok = exe("DELETE FROM generacion_programada WHERE id=%s", (int(reg_p["id"]),))
                if ok: st.success("Eliminado."); st.cache_data.clear(); st.rerun()


# ── Generación real manual (respaldo) ────────────────────────
st.markdown('<div class="sec">GENERACIÓN REAL · INGRESO MANUAL DE RESPALDO</div>', unsafe_allow_html=True)
st.info(
    "Los datos de generación real se importan automáticamente desde la API CEN (SIPUB) cada hora. "
    "El ingreso manual permite agregar o corregir valores cuando la adquisición falla. "
    "Si ya existe un registro para esa unidad/hora, el valor se sobreescribe.",
)

tab_r1, tab_r2 = st.tabs(["Por hora", "24 horas de una vez"])

with tab_r1:
    u_real = st.radio("Unidad", ["ANG1","ANG2","CCR1","CCR2"], key="ur", horizontal=True)
    rc2, rc3, rc4 = st.columns(3)
    with rc2: f_real = st.date_input("Fecha", value=hoy, max_value=hoy, key="fr")
    with rc3: h_real = st.number_input("Hora (1-24)", 1, 24, datetime.now().hour+1, key="hr")
    with rc4: mw_real = st.number_input("MW reales", 0.0, 400.0, step=0.5, key="mwr")
    if st.button("Guardar hora real", key="btn_hr"):
        fh_r = f"{f_real} {int(h_real)-1:02d}:00:00"
        ok = exe("""INSERT INTO generacion_real (unidad, gen_real_mw, fecha_hora, hora)
                    VALUES (%s, %s, %s, %s)
                    ON CONFLICT (unidad, fecha_hora) DO UPDATE SET gen_real_mw=EXCLUDED.gen_real_mw""",
                 (u_real, mw_real, fh_r, int(h_real)))
        if ok: st.success(f"Guardado: {LABELS[u_real]} H{h_real} → {mw_real} MW"); st.cache_data.clear(); st.rerun()

with tab_r2:
    st.caption("Selecciona unidad y fecha, luego pega los 24 valores MW separados por salto de línea (hora 1 → hora 24).")
    u_rmasa = st.radio("Unidad", ["ANG1","ANG2","CCR1","CCR2"], key="urm", horizontal=True)
    rm1, rm2 = st.columns([1, 2])
    with rm1:
        f_rmasa = st.date_input("Fecha", value=hoy, max_value=hoy, key="frm")
        st.caption("Ejemplo formato:\n280.5\n275.3\n271.0\n...")
    with rm2:
        mw_rmasa = st.text_area("24 valores MW (uno por línea):", height=220, key="mwrm",
                                placeholder="280.5\n275.3\n271.0\n268.5\n...")
    if st.button("Guardar las 24 horas reales", key="btn_rm"):
        try:
            valores_r = [float(v.strip().replace(",", ".")) for v in mw_rmasa.strip().split("\n") if v.strip()]
            if len(valores_r) != 24:
                st.error(f"Se esperan 24 valores, se ingresaron {len(valores_r)}.")
            else:
                params_r = [(u_rmasa, mw, f"{f_rmasa} {h-1:02d}:00:00", h) for h, mw in enumerate(valores_r, 1)]
                ok = exe_many("""INSERT INTO generacion_real (unidad, gen_real_mw, fecha_hora, hora)
                                 VALUES (%s, %s, %s, %s)
                                 ON CONFLICT (unidad, fecha_hora) DO UPDATE SET gen_real_mw=EXCLUDED.gen_real_mw""",
                              params_r)
                if ok: st.success(f"24 horas guardadas: {LABELS[u_rmasa]} · {f_rmasa}"); st.cache_data.clear(); st.rerun()
        except ValueError:
            st.error("Formato inválido. Solo números, uno por línea.")

df_rv2 = load_real(s, e)
if not df_rv2.empty:
    show_real_table = st.checkbox("Ver tabla de datos reales", value=False, key="show_real_tbl")
    if show_real_table:
        df_rshow = df_rv2.copy()
        df_rshow["fecha_hora"] = pd.to_datetime(df_rshow["fecha_hora"]).dt.strftime("%Y-%m-%d %H:%M")
        st.dataframe(df_rshow[["unidad","fecha_hora","hora","gen_real_mw"]].rename(
            columns={"gen_real_mw":"MW Real"}), use_container_width=True, hide_index=True)

    st.markdown("---")
    st.markdown("**Modificar o eliminar registro real:**")
    df_real_crud = qry(
        "SELECT id, unidad, gen_real_mw, fecha_hora, hora FROM generacion_real "
        "WHERE fecha_hora::date BETWEEN %s AND %s ORDER BY unidad, fecha_hora", (s, e))
    if not df_real_crud.empty:
        df_real_crud["fecha_hora"] = pd.to_datetime(df_real_crud["fecha_hora"])
        opc_r = df_real_crud.apply(
            lambda r: f"[{r['unidad']}] {r['fecha_hora'].strftime('%d/%m %H:%M')} — {r['gen_real_mw']:.1f} MW",
            axis=1).tolist()
        idx_r = st.selectbox("Seleccionar registro real", range(len(opc_r)),
                             format_func=lambda i: opc_r[i], key="sel_real")
        reg_r = df_real_crud.iloc[idx_r]
        col_er, col_dr = st.columns([2, 1])
        with col_er:
            new_mw_r = st.number_input("Nuevo valor MW:", value=float(reg_r["gen_real_mw"]),
                                       min_value=0.0, max_value=400.0, step=0.5, key="edit_mw_r")
            if st.button("Actualizar MW real", key="upd_real"):
                ok = exe("UPDATE generacion_real SET gen_real_mw=%s WHERE id=%s",
                         (new_mw_r, int(reg_r["id"])))
                if ok: st.success("Actualizado."); st.cache_data.clear(); st.rerun()
        with col_dr:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Eliminar", key="del_real", type="primary"):
                ok = exe("DELETE FROM generacion_real WHERE id=%s", (int(reg_r["id"]),))
                if ok: st.success("Eliminado."); st.cache_data.clear(); st.rerun()


# ── Datos horarios ────────────────────────────────────────────
st.markdown('<div class="sec">DATOS HORARIOS</div>', unsafe_allow_html=True)
show_table = st.checkbox("Ver tabla completa", value=False)
if show_table:
    df_pv = df_r.pivot_table(index="fecha_hora",columns="unidad",values="gen_real_mw",aggfunc="mean").reset_index()
    if not df_c.empty:
        df_pv = df_pv.merge(df_c[["fecha_hora","cmg_usd_mwh"]].rename(columns={"cmg_usd_mwh":"CMG (USD/MWh)"}),on="fecha_hora",how="left")
    df_pv["fecha_hora"] = pd.to_datetime(df_pv["fecha_hora"]).dt.strftime("%Y-%m-%d %H:%M")
    st.dataframe(df_pv, use_container_width=True, hide_index=True)
    st.download_button("Descargar CSV", df_pv.to_csv(index=False).encode(), f"complejo-termico-mejillones_{s}.csv","text/csv")


# ── Bitácora ──────────────────────────────────────────────────
st.markdown('<div class="sec">BITÁCORA DE NOVEDADES OPERACIONALES</div>', unsafe_allow_html=True)
tab_b1,tab_b2 = st.tabs(["Ver registros","Nueva novedad"])
with tab_b1:
    fu   = st.radio("Filtrar", "Todas ANG1 ANG2 CCR1 CCR2".split(), horizontal=True, label_visibility="collapsed")
    df_b = load_bit(s,e,fu)
    if not df_b.empty:
        df_b2 = df_b.copy()
        df_b2["fecha"] = df_b2["fecha"].astype(str)
        df_b2["hora"]  = df_b2["hora"].astype(str).str[:5]
        st.dataframe(df_b2.drop(columns=["id"],errors="ignore"),use_container_width=True,hide_index=True)

        st.markdown("---")
        st.markdown("**Modificar o eliminar registro:**")
        opciones_b = df_b.apply(lambda r: f"[{r['unidad']}] {str(r['fecha'])} {str(r['hora'])[:5]} — {str(r['comentario'])[:50]}", axis=1).tolist()
        idx_b = st.selectbox("Seleccionar registro", range(len(opciones_b)),
                             format_func=lambda i: opciones_b[i], key="sel_bit")
        reg_b = df_b.iloc[idx_b]
        col_edit_b, col_del_b = st.columns([3,1])
        with col_edit_b:
            new_com_b = st.text_area("Editar comentario:", value=str(reg_b.get("comentario","")), height=80, key="edit_com_b")
            if st.button("Actualizar novedad", key="upd_bit"):
                ok = exe("UPDATE bitacora SET comentario=%s WHERE id=%s", (new_com_b.strip(), int(reg_b["id"])))
                if ok: st.success("Actualizado."); st.cache_data.clear(); st.rerun()
        with col_del_b:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Eliminar", key="del_bit", type="primary"):
                ok = exe("DELETE FROM bitacora WHERE id=%s", (int(reg_b["id"]),))
                if ok: st.success("Eliminado."); st.cache_data.clear(); st.rerun()
    else:
        st.info("Sin novedades para el período seleccionado.")
with tab_b2:
    ub = st.radio("Unidad", ["ANG1","ANG2","CCR1","CCR2"], key="ub", horizontal=True)
    b2,b3,b4 = st.columns([1,1,1])
    with b2: ab = st.text_input("Autor / Turno",key="ab")
    with b3: fb = st.date_input("Fecha del evento",value=date.today(),key="fb")
    with b4: hb = st.time_input("Hora evento",value=datetime.now().time(),step=60,key="hb")
    cb = st.text_area("Comentario",height=90,placeholder="Descripción de la novedad operacional...",key="cb")
    if st.button("Guardar novedad",type="primary"):
        if ab.strip() and cb.strip():
            ok = exe("INSERT INTO bitacora (unidad,autor,comentario,fecha,hora) VALUES (%s,%s,%s,%s,%s)",
                     (ub,ab.strip(),cb.strip(),str(fb),str(hb)))
            if ok: st.success("Novedad guardada."); st.cache_data.clear(); st.rerun()
        else: st.warning("Completa autor y comentario.")

# ── Footer ───────────────────────────────────────────────────────
st.markdown("""
<div style="margin-top:3rem;padding-top:1rem;border-top:1px solid #E2E8F0;
            text-align:center;font-size:0.78rem;color:#94A3B8;">
    Dashboard creado por <strong style="color:#64748B;">Erick Herrera</strong>
</div>
""", unsafe_allow_html=True)