"""
utils/reports.py — Reportes EJECUTIVOS (gerencia) PDF y PPT del dashboard CTM Mejillones.

Rediseño 2026-08-01: se pasó de un reporte extenso (una página/slide por unidad,
tablas largas de SSCC, limitaciones y bitácora) a un informe ejecutivo BREVE:

  PDF → 2 páginas   (1. Resumen ejecutivo + KPIs + tabla por unidad + destacados
                     2. Gráficos consolidados de generación y CMG + eventos)
  PPT → 4 diapositivas (portada · resumen ejecutivo · desempeño · eventos y riesgos)

El detalle operacional fino sigue disponible en el dashboard; el reporte responde
solo las preguntas de gerencia: cuánto generamos, cuánto valió, qué tan bien
seguimos el programa y qué eventos lo explican.

Paleta corporativa AES (verde→teal→cyan→azul→violeta). Gráficos matplotlib
compartidos por PDF y PPT.
"""
import io
from datetime import datetime

import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

# ── Paleta corporativa AES ────────────────────────────────────────────────────
AES_SPECTRUM = ["#22A95B", "#12B2A0", "#1FB6E5", "#3D53E8", "#7C4DE0"]  # gradiente
C_DARK   = "#0F172A"
C_AZUL   = "#2A38C9"
C_CMG    = "#7C4DE0"
C_GRAY   = "#475569"
C_LGRAY  = "#94A3B8"
C_LINE   = "#E2E8F0"
C_BG     = "#F8FAFC"
C_AMBER  = "#D97706"
C_ROJO   = "#DC2626"
C_VERDE  = "#22A95B"

_UNIT_COLORS = {
    "ANG1": {"real": "#7C4DE0", "prog": "#C9B5F2"},
    "ANG2": {"real": "#3D53E8", "prog": "#AEB8F5"},
    "CCR1": {"real": "#1FB6E5", "prog": "#A6E4F6"},
    "CCR2": {"real": "#22A95B", "prog": "#A4E0BC"},
}
_UNIT_NAMES = {"ANG1": "Angamos 1", "ANG2": "Angamos 2", "CCR1": "Cochrane 1", "CCR2": "Cochrane 2"}
_UNIDADES = ["ANG1", "ANG2", "CCR1", "CCR2"]
_PMAX = {"ANG1": 277, "ANG2": 280, "CCR1": 276, "CCR2": 276}

# Potencia real < 5 MW = unidad detenida (regla 23 del proyecto: el SCADA rara
# vez marca 0.0 exacto en un trip / desconexión / mantención).
UMBRAL_TRIP = 5.0


# ══════════════════════════════════════════════════════════════
# MÉTRICAS EJECUTIVAS
# ══════════════════════════════════════════════════════════════
def _metricas_unidad(df_real, df_prog, df_cmg, unidad):
    """KPIs de gerencia de una unidad: energía, factor de planta, disponibilidad,
    adherencia al programa e ingreso estimado (energía × CMG horario)."""
    vacio = {"energia_mwh": 0.0, "prom": None, "fp": None, "disp": None,
             "horas_det": 0, "desv": None, "mae": None, "ingreso": None, "horas": 0}
    if df_real.empty:
        return vacio
    df_u = df_real[df_real["unidad"] == unidad].sort_values("fecha_hora")
    if df_u.empty:
        return vacio

    gen = pd.to_numeric(df_u["gen_real_mw"], errors="coerce").fillna(0.0)
    horas = len(df_u)
    # Datos horarios → 1 MW sostenido durante 1 h = 1 MWh.
    energia = float(gen.sum())
    prom = float(gen.mean())
    detenida = gen < UMBRAL_TRIP

    m = {
        "energia_mwh": energia,
        "prom": prom,
        "fp": prom / _PMAX[unidad] * 100.0,
        "disp": float((~detenida).mean() * 100.0),
        "horas_det": int(detenida.sum()),
        "horas": horas,
        "desv": None, "mae": None, "ingreso": None,
    }

    # Adherencia al programa (sesgo medio y error absoluto medio)
    if not df_prog.empty:
        df_up = df_prog[df_prog["unidad"] == unidad]
        if not df_up.empty:
            df_m = pd.merge_asof(
                df_u[["fecha_hora", "gen_real_mw"]].sort_values("fecha_hora"),
                df_up[["fecha_hora", "gen_programada_mw"]].sort_values("fecha_hora"),
                on="fecha_hora", direction="nearest", tolerance=pd.Timedelta("1h")).dropna()
            if not df_m.empty:
                dif = df_m["gen_real_mw"] - df_m["gen_programada_mw"]
                m["desv"] = float(dif.mean())
                m["mae"] = float(dif.abs().mean())

    # Ingreso estimado = Σ (MWh horario × CMG de esa hora)
    if df_cmg is not None and not df_cmg.empty:
        df_i = pd.merge(
            df_u[["fecha_hora", "gen_real_mw"]],
            df_cmg[["fecha_hora", "cmg_usd_mwh"]], on="fecha_hora", how="inner")
        if not df_i.empty:
            m["ingreso"] = float((df_i["gen_real_mw"] * df_i["cmg_usd_mwh"]).sum())
    return m


def _metricas_complejo(df_real, df_prog, df_cmg):
    """Agrega las 4 unidades + estadísticos de CMG del período."""
    por_u = {u: _metricas_unidad(df_real, df_prog, df_cmg, u) for u in _UNIDADES}
    energia = sum(m["energia_mwh"] for m in por_u.values())
    ingresos = [m["ingreso"] for m in por_u.values() if m["ingreso"] is not None]
    fps = [m["fp"] for m in por_u.values() if m["fp"] is not None]
    disps = [m["disp"] for m in por_u.values() if m["disp"] is not None]
    maes = [m["mae"] for m in por_u.values() if m["mae"] is not None]

    cmg = {}
    if df_cmg is not None and not df_cmg.empty:
        s = pd.to_numeric(df_cmg["cmg_usd_mwh"], errors="coerce").dropna()
        if not s.empty:
            cmg = {"prom": float(s.mean()), "min": float(s.min()), "max": float(s.max()),
                   "horas_cero": int((s <= 0.5).sum()), "horas": int(len(s))}

    return {
        "por_unidad": por_u,
        "energia_mwh": energia,
        "energia_gwh": energia / 1000.0,
        "ingreso": sum(ingresos) if ingresos else None,
        "precio_medio": (sum(ingresos) / energia) if ingresos and energia > 0 else None,
        "fp": (sum(fps) / len(fps)) if fps else None,
        "disp": (sum(disps) / len(disps)) if disps else None,
        "mae": (sum(maes) / len(maes)) if maes else None,
        "horas_det": sum(m["horas_det"] for m in por_u.values()),
        "cmg": cmg,
    }


def _destacados(kpi, df_sscc, df_lim, nodo_label="Crucero 220 kV"):
    """Bullets narrativos para gerencia, derivados de los KPIs y los eventos."""
    out = []
    por_u = kpi["por_unidad"]

    if kpi["energia_gwh"] > 0:
        txt = f"El complejo generó {kpi['energia_gwh']:.1f} GWh"
        if kpi["ingreso"] is not None:
            txt += (f", con un ingreso estimado de {kpi['ingreso']/1000:,.0f} kUSD"
                    f" (precio medio {kpi['precio_medio']:.1f} USD/MWh)")
        if kpi["fp"] is not None:
            txt += f" y un factor de planta medio de {kpi['fp']:.0f}%"
        out.append(txt + ".")

    # Unidad líder y unidad rezagada por energía
    ranking = sorted(((u, m["energia_mwh"]) for u, m in por_u.items() if m["energia_mwh"] > 0),
                     key=lambda x: -x[1])
    if len(ranking) >= 2:
        top, bot = ranking[0], ranking[-1]
        out.append(f"{_UNIT_NAMES[top[0]]} lideró el despacho con {top[1]/1000:.1f} GWh; "
                   f"{_UNIT_NAMES[bot[0]]} fue la de menor aporte con {bot[1]/1000:.1f} GWh.")

    if kpi["cmg"]:
        c = kpi["cmg"]
        txt = (f"CMG {nodo_label}: promedio {c['prom']:.1f} USD/MWh "
               f"(rango {c['min']:.1f} – {c['max']:.1f}).")
        if c["horas_cero"] > 0:
            pct = c["horas_cero"] / max(c["horas"], 1) * 100
            txt += (f" Se registraron {c['horas_cero']} horas con CMG en cero "
                    f"({pct:.0f}% del período), señal de desacople de la zona.")
        out.append(txt)

    if kpi["mae"] is not None:
        calidad = "alta" if kpi["mae"] < 15 else ("media" if kpi["mae"] < 35 else "baja")
        out.append(f"Adherencia al programa: error absoluto medio de {kpi['mae']:.1f} MW "
                   f"por unidad (precisión {calidad}).")

    if kpi["horas_det"] > 0:
        det = [f"{_UNIT_NAMES[u]} ({m['horas_det']} h)"
               for u, m in por_u.items() if m["horas_det"] > 0]
        out.append(f"Horas con unidad detenida (< {UMBRAL_TRIP:.0f} MW): "
                   f"{kpi['horas_det']} en total — " + ", ".join(det) + ".")
    else:
        out.append("Sin horas de detención: las cuatro unidades operaron sobre el umbral "
                   f"de {UMBRAL_TRIP:.0f} MW durante todo el período.")

    if df_lim is not None and not df_lim.empty:
        # Se informa lo que estuvo VIGENTE en el período, no el conteo de
        # «pendientes»: el CEN nunca cierra los registros, así que ese número
        # arrastraba limitaciones de meses atrás ya superadas (utils/eventos.py).
        n_act = int((df_lim["_estado"] == "activa").sum()) if "_estado" in df_lim.columns else 0
        n_sscc = int(df_lim["afecta_sscc"].fillna(False).sum()) if "afecta_sscc" in df_lim.columns else 0
        if n_act:
            out.append(f"Limitaciones de transmisión: {len(df_lim)} vigentes en el período, "
                       f"{n_act} aún abiertas al cierre y {n_sscc} con impacto en SSCC.")
        else:
            out.append(f"Limitaciones de transmisión: {len(df_lim)} afectaron el período y "
                       f"ninguna sigue vigente al cierre"
                       + (f"; {n_sscc} con impacto en SSCC." if n_sscc else "."))

    if df_sscc is not None and not df_sscc.empty:
        out.append(f"Servicios complementarios: {len(df_sscc)} instrucciones recibidas del CEN.")

    return out


# ══════════════════════════════════════════════════════════════
# HELPERS — gráficos matplotlib
# ══════════════════════════════════════════════════════════════
def _estilo_ejes(ax, ylabel):
    ax.set_ylabel(ylabel, fontsize=9, color=C_GRAY)
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%a %d/%m"))
    ax.xaxis.set_major_locator(mdates.DayLocator())
    plt.setp(ax.xaxis.get_majorticklabels(), fontsize=8, rotation=0, color=C_GRAY)
    ax.tick_params(axis="y", labelsize=8, colors=C_GRAY)
    ax.set_facecolor("#FCFDFF")
    ax.grid(axis="y", color=C_LINE, linewidth=0.7)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    for s in ("left", "bottom"):
        ax.spines[s].set_color(C_LINE)


def _fig_generacion_complejo(df_real, figsize=(14, 4.2)):
    """Una sola figura con las 4 unidades — reemplaza los 4 gráficos por unidad."""
    fig, ax = plt.subplots(figsize=figsize)
    if df_real.empty:
        ax.text(0.5, 0.5, "Sin datos de generación", ha="center", va="center",
                transform=ax.transAxes, color=C_LGRAY)
    else:
        for u in _UNIDADES:
            df_u = df_real[df_real["unidad"] == u].sort_values("fecha_hora")
            if df_u.empty:
                continue
            ax.plot(df_u["fecha_hora"], df_u["gen_real_mw"], color=_UNIT_COLORS[u]["real"],
                    linewidth=1.9, label=_UNIT_NAMES[u])
        ax.legend(loc="upper center", bbox_to_anchor=(0.5, 1.13), ncol=4,
                  fontsize=8, frameon=False)
    ax.set_ylim(bottom=0)
    _estilo_ejes(ax, "MW")
    fig.patch.set_facecolor("white")
    fig.tight_layout(pad=0.4)
    return fig


def _fig_cmg(df_cmg, figsize=(14, 3)):
    fig, ax = plt.subplots(figsize=figsize)
    if df_cmg is None or df_cmg.empty:
        ax.text(0.5, 0.5, "Sin datos CMG", ha="center", va="center",
                transform=ax.transAxes, color=C_LGRAY)
    else:
        ax.plot(df_cmg["fecha_hora"], df_cmg["cmg_usd_mwh"], color=C_CMG, linewidth=2.0)
        ax.fill_between(df_cmg["fecha_hora"], df_cmg["cmg_usd_mwh"], alpha=0.10, color=C_CMG)
        prom = df_cmg["cmg_usd_mwh"].mean()
        ax.axhline(prom, color=C_LGRAY, linewidth=1, linestyle=":", zorder=1)
        ax.annotate(f"Prom {prom:.1f}", xy=(0.995, prom), xycoords=("axes fraction", "data"),
                    ha="right", va="bottom", fontsize=8, color=C_GRAY)
    _estilo_ejes(ax, "USD/MWh")
    fig.patch.set_facecolor("white")
    fig.tight_layout(pad=0.4)
    return fig


def _fig_energia_barras(kpi, figsize=(6.5, 3.0)):
    """Energía por unidad (GWh) — lectura inmediata del aporte relativo."""
    fig, ax = plt.subplots(figsize=figsize)
    vals = [kpi["por_unidad"][u]["energia_mwh"] / 1000.0 for u in _UNIDADES]
    if sum(vals) <= 0:
        ax.text(0.5, 0.5, "Sin datos", ha="center", va="center",
                transform=ax.transAxes, color=C_LGRAY)
    else:
        ax.bar([_UNIT_NAMES[u] for u in _UNIDADES], vals,
               color=[_UNIT_COLORS[u]["real"] for u in _UNIDADES], width=0.6)
        for i, v in enumerate(vals):
            ax.text(i, v, f"{v:.1f}", ha="center", va="bottom", fontsize=9,
                    color=C_DARK, fontweight="bold")
        ax.set_ylim(0, max(vals) * 1.18)
    ax.set_ylabel("GWh", fontsize=9, color=C_GRAY)
    ax.tick_params(axis="both", labelsize=8, colors=C_GRAY)
    ax.set_facecolor("#FCFDFF")
    ax.grid(axis="y", color=C_LINE, linewidth=0.7)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    for s in ("left", "bottom"):
        ax.spines[s].set_color(C_LINE)
    fig.patch.set_facecolor("white")
    fig.tight_layout(pad=0.4)
    return fig


def _fig_to_bytes(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _periodo(start_str, end_str):
    try:
        dt_s = datetime.strptime(start_str, "%Y-%m-%d")
        dt_e = datetime.strptime(end_str, "%Y-%m-%d")
        return (dt_s.isocalendar()[1], dt_s.strftime("%d/%m/%Y"), dt_e.strftime("%d/%m/%Y"),
                (dt_e - dt_s).days + 1)
    except Exception:
        return "—", start_str, end_str, 0


def _f(v, fmt="{:.1f}", suf=""):
    return "—" if v is None or pd.isna(v) else (fmt.format(v) + suf)


def _filas_unidad(kpi):
    """Filas de la tabla ejecutiva por unidad (mismas en PDF y PPT)."""
    filas = []
    for u in _UNIDADES:
        m = kpi["por_unidad"][u]
        filas.append([
            _UNIT_NAMES[u],
            _f(m["energia_mwh"] / 1000.0 if m["energia_mwh"] else 0.0, "{:.2f}"),
            _f(m["prom"], "{:.0f}"),
            _f(m["fp"], "{:.0f}", "%"),
            _f(m["disp"], "{:.0f}", "%"),
            _f(m["desv"], "{:+.1f}"),
            _f(m["ingreso"] / 1000.0 if m["ingreso"] is not None else None, "{:,.0f}"),
        ])
    return filas


_ENCABEZADOS = ["Unidad", "Energía (GWh)", "Prom (MW)", "Factor planta",
                "Disponibilidad", "Desv. vs prog. (MW)", "Ingreso (kUSD)"]


def _fuente_programa(df_prog):
    """Texto corto con el ORIGEN del programa usado como referencia.

    Importa dejarlo escrito: el CEN a veces deja de emitir el PID y la serie se
    rellena con el PCP (visto desde 2026-07-08). Sin esta nota, una desviación
    'vs programa' se leería como si siempre comparara contra lo mismo.
    """
    if df_prog is None or df_prog.empty or "fuente" not in df_prog.columns:
        return "Sin programa de referencia en el período."
    nombres = {"CEN_PID": "PID (intra-día)", "CEN_PCP": "PCP (día-ante)", "MANUAL": "ingreso manual"}
    cuenta = df_prog["fuente"].value_counts()
    total = int(cuenta.sum())
    partes = [f"{nombres.get(f, f)} {n / total * 100:.0f}%" for f, n in cuenta.items()]
    txt = "Programa de referencia: " + " · ".join(partes) + "."
    if "CEN_PID" not in cuenta.index:
        txt += " El CEN no emitió PID en el período; se usa el PCP como programa."
    return txt


def _eventos_resumen(df_sscc, df_lim):
    """Resumen compacto de eventos por unidad (conteos, no listados)."""
    filas = []
    for u in _UNIDADES:
        n_sscc = int((df_sscc["unidad"] == u).sum()) if df_sscc is not None and not df_sscc.empty else 0
        if df_lim is not None and not df_lim.empty and "_unidad" in df_lim.columns:
            df_lu = df_lim[df_lim["_unidad"] == u]
            n_lim = len(df_lu)
            # Abiertas al cierre = ventana aún vigente, no «status pendiente».
            n_pend = (int((df_lu["_estado"] == "activa").sum())
                      if n_lim and "_estado" in df_lu.columns else 0)
        else:
            n_lim = n_pend = 0
        filas.append([_UNIT_NAMES[u], str(n_sscc), str(n_lim), str(n_pend)])
    return filas


# ══════════════════════════════════════════════════════════════
# GENERADOR PDF — 2 páginas
# ══════════════════════════════════════════════════════════════
def generar_pdf(df_real, df_prog, df_cmg, start_str, end_str, df_sscc=None, df_lim=None,
                nodo_label="Crucero 220 kV"):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor, white
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Image as RLImage, HRFlowable, PageBreak, Table, TableStyle)
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT

    if df_sscc is None: df_sscc = pd.DataFrame()
    if df_lim  is None: df_lim  = pd.DataFrame()

    kpi = _metricas_complejo(df_real, df_prog, df_cmg)
    semana, fmt_s, fmt_e, n_dias = _periodo(start_str, end_str)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=1.6*cm, rightMargin=1.6*cm,
                            topMargin=1.3*cm, bottomMargin=1.2*cm,
                            title="CTM Mejillones — Reporte Ejecutivo")

    H_DARK, H_AZUL = HexColor(C_DARK), HexColor(C_AZUL)
    H_GRAY, H_LGRAY, H_LINE = HexColor(C_GRAY), HexColor(C_LGRAY), HexColor(C_LINE)

    def _sty(name, size=10, color=None, after=4, align=TA_LEFT, bold=False, leading=None):
        return ParagraphStyle(name, fontName="Helvetica-Bold" if bold else "Helvetica",
                              fontSize=size, leading=leading or size * 1.25,
                              textColor=color or H_DARK, spaceAfter=after, alignment=align)

    sTitle  = _sty("tt", size=21, bold=True, after=2, color=H_DARK)
    sSub    = _sty("ss", size=10, color=H_GRAY, after=1)
    sH2     = _sty("h2", size=12, bold=True, after=5, color=H_AZUL)
    sBullet = _sty("bl", size=9.2, color=H_DARK, after=5, leading=13)
    sSmall  = _sty("sm", size=7.5, color=H_LGRAY, after=2)
    sCenter = _sty("ct", size=7.5, color=H_LGRAY, after=2, align=TA_CENTER)

    ANCHO = 17.8 * cm

    def _accent_bar(width=ANCHO):
        w = width / len(AES_SPECTRUM)
        t = Table([[""] * len(AES_SPECTRUM)], colWidths=[w] * len(AES_SPECTRUM),
                  rowHeights=[0.14*cm])
        sty = [("TOPPADDING", (0,0), (-1,-1), 0), ("BOTTOMPADDING", (0,0), (-1,-1), 0)]
        for i, c in enumerate(AES_SPECTRUM):
            sty.append(("BACKGROUND", (i,0), (i,0), HexColor(c)))
        t.setStyle(TableStyle(sty))
        return t

    story = []

    # ── ENCABEZADO ─────────────────────────────────────────────
    story += [
        _accent_bar(), Spacer(1, 0.32*cm),
        Paragraph("Complejo Térmico Mejillones", sTitle),
        Paragraph("Reporte ejecutivo de operación  ·  Angamos 1-2 · Cochrane 1-2", sSub),
        Paragraph(f"Período: {fmt_s} al {fmt_e}  ·  Semana {semana}  ·  {n_dias} días", sSub),
        Spacer(1, 0.45*cm),
    ]

    # ── FILA DE KPIs DEL COMPLEJO ──────────────────────────────
    cmg_prom = kpi["cmg"].get("prom") if kpi["cmg"] else None
    kpi_cells = [
        ("Energía generada", _f(kpi["energia_gwh"], "{:.2f}", " GWh")),
        ("Ingreso estimado", _f(kpi["ingreso"] / 1000.0 if kpi["ingreso"] is not None else None,
                                "{:,.0f}", " kUSD")),
        ("Factor de planta", _f(kpi["fp"], "{:.0f}", "%")),
        ("Disponibilidad", _f(kpi["disp"], "{:.0f}", "%")),
        ("CMG promedio", _f(cmg_prom, "{:.1f}", " USD/MWh")),
    ]
    tk = Table([[c[0] for c in kpi_cells], [c[1] for c in kpi_cells]],
               colWidths=[ANCHO / 5] * 5)
    tk.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), HexColor(C_BG)),
        ("FONTNAME", (0,0), (-1,0), "Helvetica"), ("FONTSIZE", (0,0), (-1,0), 7.5),
        ("TEXTCOLOR", (0,0), (-1,0), H_GRAY),
        ("FONTNAME", (0,1), (-1,1), "Helvetica-Bold"), ("FONTSIZE", (0,1), (-1,1), 13),
        ("TEXTCOLOR", (0,1), (-1,1), H_AZUL),
        ("ALIGN", (0,0), (-1,-1), "CENTER"), ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("LINEAFTER", (0,0), (-2,-1), 0.5, white),
        ("TOPPADDING", (0,0), (-1,0), 7), ("BOTTOMPADDING", (0,0), (-1,0), 1),
        ("TOPPADDING", (0,1), (-1,1), 0), ("BOTTOMPADDING", (0,1), (-1,1), 8),
    ]))
    story += [tk, Spacer(1, 0.5*cm)]

    # ── DESTACADOS ─────────────────────────────────────────────
    story += [Paragraph("Puntos destacados", sH2)]
    for b in _destacados(kpi, df_sscc, df_lim, nodo_label):
        story.append(Paragraph(f"•&nbsp;&nbsp;{b}", sBullet))
    story += [Spacer(1, 0.35*cm)]

    # ── TABLA POR UNIDAD ───────────────────────────────────────
    story += [Paragraph("Desempeño por unidad", sH2)]
    rows = [_ENCABEZADOS] + _filas_unidad(kpi)
    col_w = [3.0*cm, 2.5*cm, 2.2*cm, 2.4*cm, 2.6*cm, 2.9*cm, 2.2*cm]
    tbl = Table(rows, colWidths=col_w)
    sty = [
        ("BACKGROUND", (0,0), (-1,0), H_AZUL),
        ("TEXTCOLOR",  (0,0), (-1,0), white),
        ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTNAME",   (0,1), (0,-1), "Helvetica-Bold"),
        ("FONTSIZE",   (0,0), (-1,-1), 8.5),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [white, HexColor(C_BG)]),
        ("GRID",  (0,0), (-1,-1), 0.4, H_LINE),
        ("ALIGN", (1,0), (-1,-1), "CENTER"), ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5), ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]
    for i, u in enumerate(_UNIDADES, 1):
        sty.append(("TEXTCOLOR", (0,i), (0,i), HexColor(_UNIT_COLORS[u]["real"])))
    tbl.setStyle(TableStyle(sty))
    story += [tbl, Spacer(1, 0.15*cm),
              Paragraph("Energía y potencia media sobre datos horarios del CEN. Disponibilidad = "
                        f"horas con potencia ≥ {UMBRAL_TRIP:.0f} MW. Ingreso estimado = Σ (MWh × CMG "
                        f"{nodo_label} de la hora); es una referencia de mercado, no una liquidación.",
                        sSmall),
              Paragraph(_fuente_programa(df_prog), sSmall), Spacer(1, 0.5*cm)]

    # ── Aporte por unidad + eventos, lado a lado (cierra la página 1) ──
    img_barras = RLImage(_fig_to_bytes(_fig_energia_barras(kpi, figsize=(6.4, 3.0))),
                         width=8.6*cm, height=4.0*cm)
    ev_rows = [["Unidad", "SSCC", "Limitac.", "Vigentes"]] + _eventos_resumen(df_sscc, df_lim)
    t_ev = Table(ev_rows, colWidths=[3.0*cm, 1.9*cm, 2.0*cm, 2.2*cm])
    t_ev.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), HexColor(C_AMBER)),
        ("TEXTCOLOR", (0,0), (-1,0), white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTNAME", (0,1), (0,-1), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8.5),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [white, HexColor(C_BG)]),
        ("GRID", (0,0), (-1,-1), 0.4, H_LINE),
        ("ALIGN", (1,0), (-1,-1), "CENTER"), ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5), ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))
    bloque = Table([[img_barras,
                     [Paragraph("Eventos del período", sH2), t_ev]]],
                   colWidths=[8.9*cm, 8.9*cm])
    bloque.setStyle(TableStyle([("VALIGN", (0,0), (-1,-1), "TOP"),
                                ("LEFTPADDING", (0,0), (-1,-1), 0),
                                ("RIGHTPADDING", (0,0), (-1,-1), 0)]))
    story += [Paragraph("Aporte por unidad y eventos", sH2), bloque, PageBreak()]

    # ── PÁGINA 2: SERIES DE TIEMPO ─────────────────────────────
    story += [_accent_bar(), Spacer(1, 0.3*cm),
              Paragraph("Generación del complejo y costo marginal", sH2)]
    story.append(RLImage(_fig_to_bytes(_fig_generacion_complejo(df_real, figsize=(14, 6.4))),
                         width=ANCHO, height=8.2*cm))
    story.append(Spacer(1, 0.4*cm))
    story.append(RLImage(_fig_to_bytes(_fig_cmg(df_cmg, figsize=(14, 4.4))),
                         width=ANCHO, height=5.7*cm))
    if kpi["cmg"]:
        c = kpi["cmg"]
        story.append(Paragraph(
            f"CMG {nodo_label} — prom {c['prom']:.1f} · mín {c['min']:.1f} · máx {c['max']:.1f} USD/MWh"
            f" · {c['horas_cero']} horas en cero.", sSmall))
    story.append(Spacer(1, 0.5*cm))

    story += [
        HRFlowable(width="100%", thickness=0.5, color=H_LINE),
        Paragraph(f"Generado {datetime.now().strftime('%d/%m/%Y %H:%M')}  ·  "
                  "Fuente: API CEN (SIP / Operaciones)  ·  Dashboard CTM Mejillones", sCenter),
    ]

    doc.build(story)
    buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════
# GENERADOR PPT — 4 diapositivas
# ══════════════════════════════════════════════════════════════
def generar_ppt(df_real, df_prog, df_cmg, start_str, end_str, df_sscc=None, df_lim=None,
                nodo_label="Crucero 220 kV"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if df_sscc is None: df_sscc = pd.DataFrame()
    if df_lim  is None: df_lim  = pd.DataFrame()

    kpi = _metricas_complejo(df_real, df_prog, df_cmg)
    semana, fmt_s, fmt_e, n_dias = _periodo(start_str, end_str)

    def _rgb(hexs):
        h = hexs.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    RGB_DARK, RGB_AZUL = _rgb(C_DARK), _rgb(C_AZUL)
    RGB_GRAY, RGB_LGRAY = _rgb(C_GRAY), _rgb(C_LGRAY)
    RGB_WHITE, RGB_BG   = RGBColor(0xFF, 0xFF, 0xFF), _rgb(C_BG)
    RGB_AMBER, RGB_CMG  = _rgb(C_AMBER), _rgb(C_CMG)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    blank = prs.slide_layouts[6]

    def _slide(bg=RGB_WHITE):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb = bg
        return sl

    def _txb(sl, text, l, t, w, h, size=18, bold=False, color=RGB_DARK, align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
        return tb

    def _rect(sl, l, t, w, h, color):
        shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        return shp

    def _accent(sl, l, t, w, h=0.09):
        seg = w / len(AES_SPECTRUM)
        for i, c in enumerate(AES_SPECTRUM):
            _rect(sl, l + i*seg, t, seg, h, _rgb(c))

    def _cabecera(sl, titulo, color):
        _rect(sl, 0, 0, 13.33, 0.62, color)
        _accent(sl, 0, 0.62, 13.33, 0.06)
        _txb(sl, titulo, 0.35, 0.1, 12, 0.44, size=20, bold=True, color=RGB_WHITE)

    def _img(sl, fig, l, t, w, h):
        sl.shapes.add_picture(_fig_to_bytes(fig, dpi=150), Inches(l), Inches(t), Inches(w), Inches(h))

    def _tabla(sl, encabezados, filas, l, t, w, h, colores_1a_col=None, cab_rgb=RGB_AZUL):
        shape = sl.shapes.add_table(len(filas) + 1, len(encabezados),
                                    Inches(l), Inches(t), Inches(w), Inches(h))
        tabla = shape.table
        for j, txt in enumerate(encabezados):
            cel = tabla.cell(0, j)
            cel.text = txt
            cel.fill.solid(); cel.fill.fore_color.rgb = cab_rgb
            p = cel.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(11); p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = RGB_WHITE
        for i, fila in enumerate(filas, 1):
            for j, txt in enumerate(fila):
                cel = tabla.cell(i, j)
                cel.text = str(txt)
                cel.fill.solid()
                cel.fill.fore_color.rgb = RGB_WHITE if i % 2 else RGB_BG
                p = cel.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                r = p.runs[0]; r.font.size = Pt(11); r.font.bold = (j == 0)
                r.font.color.rgb = (colores_1a_col[i-1] if (j == 0 and colores_1a_col)
                                    else RGB_DARK)
        return tabla

    unit_rgb = [_rgb(_UNIT_COLORS[u]["real"]) for u in _UNIDADES]

    # ── SLIDE 1: PORTADA ───────────────────────────────────────
    sl = _slide(RGB_DARK)
    _accent(sl, 0.6, 2.3, 5.0, 0.12)
    _txb(sl, "Complejo Térmico Mejillones", 0.6, 2.6, 12, 1.0, size=40, bold=True, color=RGB_WHITE)
    _txb(sl, "Reporte ejecutivo de operación", 0.6, 3.7, 12, 0.6, size=18, color=RGB_LGRAY)
    _txb(sl, f"{fmt_s} — {fmt_e}  ·  Semana {semana}  ·  {n_dias} días",
         0.6, 4.4, 12, 0.5, size=16, color=_rgb("#8FA0FF"))
    _txb(sl, f"Generado {datetime.now().strftime('%d/%m/%Y %H:%M')}  ·  Fuente: API CEN",
         0.6, 6.9, 12, 0.4, size=10, color=RGB_LGRAY)

    # ── SLIDE 2: RESUMEN EJECUTIVO ─────────────────────────────
    sl = _slide(RGB_BG)
    _cabecera(sl, "Resumen ejecutivo", RGB_AZUL)
    cmg_prom = kpi["cmg"].get("prom") if kpi["cmg"] else None
    tiles = [
        ("Energía generada", _f(kpi["energia_gwh"], "{:.2f}", " GWh")),
        ("Ingreso estimado", _f(kpi["ingreso"] / 1000.0 if kpi["ingreso"] is not None else None,
                                "{:,.0f}", " kUSD")),
        ("Factor de planta", _f(kpi["fp"], "{:.0f}", "%")),
        ("Disponibilidad", _f(kpi["disp"], "{:.0f}", "%")),
        ("CMG promedio", _f(cmg_prom, "{:.1f}", " USD/MWh")),
    ]
    for i, (lbl, val) in enumerate(tiles):
        bx = 0.35 + i * 2.58
        _rect(sl, bx, 0.9, 2.42, 1.0, RGB_WHITE)
        _txb(sl, lbl, bx + 0.06, 0.96, 2.3, 0.28, size=9, color=RGB_GRAY, align=PP_ALIGN.CENTER)
        _txb(sl, val, bx + 0.06, 1.26, 2.3, 0.5, size=16, bold=True, color=RGB_AZUL,
             align=PP_ALIGN.CENTER)

    _txb(sl, "Puntos destacados", 0.35, 2.1, 6, 0.35, size=14, bold=True, color=RGB_DARK)
    y = 2.5
    for b in _destacados(kpi, df_sscc, df_lim, nodo_label)[:6]:
        _txb(sl, f"•  {b}", 0.42, y, 12.5, 0.55, size=12, color=RGB_DARK)
        y += 0.62
        if y > 6.6:
            break

    # ── SLIDE 3: DESEMPEÑO POR UNIDAD ──────────────────────────
    sl = _slide(RGB_BG)
    _cabecera(sl, "Desempeño por unidad", RGB_AZUL)
    _tabla(sl, _ENCABEZADOS, _filas_unidad(kpi), 0.35, 0.95, 12.63, 1.9,
           colores_1a_col=unit_rgb)
    _img(sl, _fig_generacion_complejo(df_real, figsize=(12.6, 3.4)), 0.35, 3.15, 12.63, 3.4)
    _txb(sl, f"Ingreso estimado = Σ (MWh × CMG {nodo_label} de la hora). "
             f"Disponibilidad = horas con potencia ≥ {UMBRAL_TRIP:.0f} MW. "
             + _fuente_programa(df_prog),
         0.35, 6.7, 12.6, 0.6, size=9, color=RGB_GRAY)

    # ── SLIDE 4: CMG, APORTE Y EVENTOS ─────────────────────────
    sl = _slide(RGB_BG)
    _cabecera(sl, "Costo marginal, aporte y eventos", RGB_CMG)
    _img(sl, _fig_cmg(df_cmg, figsize=(12.6, 2.6)), 0.35, 0.95, 12.63, 2.6)
    if kpi["cmg"]:
        c = kpi["cmg"]
        _txb(sl, f"CMG {nodo_label}  ·  prom {c['prom']:.1f}  ·  mín {c['min']:.1f}  ·  "
                 f"máx {c['max']:.1f} USD/MWh  ·  {c['horas_cero']} horas en cero",
             0.35, 3.6, 12.6, 0.4, size=11, color=RGB_GRAY)
    _img(sl, _fig_energia_barras(kpi, figsize=(6.2, 2.6)), 0.35, 4.1, 6.1, 2.6)
    _txb(sl, "Eventos del período", 6.9, 4.05, 6, 0.35, size=14, bold=True, color=RGB_DARK)
    _tabla(sl, ["Unidad", "SSCC", "Limitaciones", "Vigentes"],
           _eventos_resumen(df_sscc, df_lim), 6.9, 4.45, 6.05, 1.9,
           colores_1a_col=unit_rgb, cab_rgb=RGB_AMBER)
    _txb(sl, f"Generado {datetime.now().strftime('%d/%m/%Y %H:%M')}  ·  Fuente: API CEN "
             "(SIP / Operaciones)", 0.35, 6.95, 12.6, 0.4, size=9, color=RGB_LGRAY)

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf.read()
